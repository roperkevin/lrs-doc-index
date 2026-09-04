"""Gate for the draft -> review deck converter (local/draft2pptx.mjs).

Converts a representative TestPlanGen draft (banner comment, WARNING
alert, verifier IMPORTANT block, task lists, TC cases with Expected
Result / Trace labels, GFM tables, the Issue Trace addendum) and
reads the .pptx back with python-pptx, asserting the deck contract:

  1. the package opens and carries the expected slide walk: title,
     at-a-glance, setup checklist, a divider + case slide per test
     section, open questions, Coverage Map, Issue Trace, closing
  2. the title slide carries the plan headline, the DRAFT pill, and
     the Overview facts; the machine banner comment reaches NO slide
  3. the at-a-glance stat tiles count positive/negative cases and
     real [VERIFY: flags (banner prose mentioning the tag without a
     colon must not inflate the count)
  4. each case slide carries its TC id chip, kind chip, steps
     verbatim (numbering preserved), and the Expected Result and
     Trace card texts
  5. [VERIFY: ...] spans surface as amber (C2701A) runs
  6. Coverage Map and Issue Trace arrive as native, editable
     PowerPoint tables with their cells intact
  7. the closing slide carries the draft's provenance line
  8. CLI contract: sibling-name default, -o for a single input,
     multi-input conversion, usage errors exit nonzero
  9. figures (v1.1, prompt v1.10's FIGURES rule): a case's
     **Figure:** line + --media <dir> yields a figure slide directly
     after the case slide — TC id chip, FIGURE tag, the story's alt
     text as the title, and the cited SVG as a native shape GROUP
     (the svg2pptx emitter; the plate stays dropped, the raw URL
     reaches no slide); without --media (or with the file missing)
     the deck converts anyway, the case slide carries the muted
     "Figure: … (not embedded)" note, and stderr names the fix; a
     --media path that is not a directory is refused up front

Needs python-pptx (review/harness/requirements.txt — the CI
full-format job installs it). Usage: python3 check_draft2pptx.py
"""
import os
import subprocess
import sys
import tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
JOB = os.path.join(HERE, "..", "draft2pptx.mjs")
REPO = os.path.join(HERE, "..", "..")

PASS, FAIL = [], []


def check(name, cond, detail=""):
    (PASS if cond else FAIL).append(name)
    mark = "ok  " if cond else "FAIL"
    print(f"  {mark} {name}" + ("" if cond else f"  <- {detail}"))


DRAFT = """<!-- machine-generated test-plan draft — TestPlanGen prompt v1.7 · local/testplangen.mjs v1.3 · provider anthropic -->
> [!WARNING]
> **DRAFT — machine-generated, unreviewed.** Generated 2026-09-04T00:00:00Z from user story doc 12 — "Route Merge". Source sidecar: <https://example/sc.md>
> Review every case and resolve all [VERIFY] items before use.

<!-- verify: 1 finding(s) — lib/draftlint.mjs, prompt v1.7 contract + grounding -->
> [!IMPORTANT]
> Draft verifier: 1 finding(s) — review these first:
> - TC-N1 carries a **Trace:** line

# Test Plan — Route Merge

## Overview

| Surface | Target release | PE |
| --- | --- | --- |
| Pro | 3.8 | Claire Wang |

Verifies measure-preserving merge of two routes in ArcGIS Pro,
covering the workflow the story enumerates.

## Setup / Prerequisites
- [ ] 1. LRS network with two mergeable routes. [VERIFY: minimum network configuration]
- [x] 2. Two Pro sessions signed in.

## Positive Tests

### TC-P1 — Merge preserves measures
**Steps:**
- [ ] 1. Run Merge Routes on route A and route B.
- [ ] 2. Inspect the measures on the merged route.

**Expected Result:** The merged route keeps the source measures unchanged.

**Trace:** "the merge must preserve measures" — story requirement.

## Negative Tests

> [!CAUTION]
> A pass below is the described denial or error — never the edit
> succeeding.

### TC-N1 — Merge denied on locked route
**Steps:**
- [ ] 1. As user B, attempt Merge Routes on a locked route.

**Expected Result:** The merge is denied with a lock conflict.

**Trace:** exemplar pattern — multi-user denial case.

## Open Questions
- [ ] [VERIFY: minimum network configuration for setup]

## Coverage Map

| # | Requirement (source) | Covered by |
| --- | --- | --- |
| 1 | "the merge must preserve measures" (requirement) | TC-P1 |
| 2 | denial on locked routes (conflict statement) | TC-N1 |

## Issue Trace

_Deterministic addendum — minted by local/testplangen.mjs from the Doc IDs and Issue Refs lists, not by the model. Cross-check against devtopia during the review pass._

| Issue | Title (Issue Refs) | Schedule status | Found via |
| --- | --- | --- | --- |
| ArcGISPro/ps-location-referencing#4855 | Route merge epic | Iteration 2 · Dev=Completed | sidecar |
"""


# a compact SlideFigures-vocabulary figure (plate + routes + ticks +
# texts) — the link shape the sweep writes and prompt v1.10 cites
FIG_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 760 320" width="760" '
    'height="320" role="img" aria-label="Slide 2 diagram">'
    '<title>Slide 2 — Routes before the merge</title>'
    '<desc>Routes R1 and R2 end to end.</desc>'
    "<style>.plate{fill:#FFFFFF;stroke:#D7DFDF;stroke-width:1}"
    ".ln{fill:none;stroke-linecap:round;stroke-linejoin:round}"
    ".route{stroke:#16302F;stroke-width:3;stroke-linecap:butt}"
    ".tick{stroke:#6E8285;stroke-width:1.15}"
    ".measure{font-size:11px;fill:#6E8285}.id{font-size:12.5px;font-weight:600}"
    ".f-cool{fill:#1B6E8C}text{font-family:'Segoe UI',Arial,sans-serif}</style>"
    '<rect class="plate" x="1" y="1" width="758" height="318" rx="6"/>'
    '<g transform="translate(20,40)">'
    '<line class="ln route" x1="40" y1="60" x2="340" y2="60"/>'
    '<line class="ln route" x1="360" y1="60" x2="660" y2="60"/>'
    '<line class="ln tick" x1="40" y1="52" x2="40" y2="68"/>'
    '<text class="measure" x="40" y="44" text-anchor="middle">0</text>'
    '<text class="id f-cool" x="190" y="84" text-anchor="middle">R1</text>'
    '<text class="id f-cool" x="510" y="84" text-anchor="middle">R2</text>'
    "</g></svg>")

FIG_URL = ("https://mock.example/sites/lrsworkspace/LRS%20Doc%20Index/media/"
           "doc12_slide2_fig1.svg")
FIG_ALT = "Routes R1 and R2 before the merge"
FIG_DRAFT = DRAFT.replace(
    '**Trace:** "the merge must preserve measures" — story requirement.',
    '**Trace:** "the merge must preserve measures" — story requirement.\n\n'
    f'**Figure:** ![{FIG_ALT}]({FIG_URL})')


def slide_text(slide):
    """Every text run on a slide, joined — shapes and table cells."""
    parts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            parts.append(sh.text_frame.text)
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


def runs_of(slide):
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    yield r


def main():
    try:
        import pptx
        from pptx.util import Emu  # noqa: F401
    except ImportError:
        print("SKIP: python-pptx not installed (pip install -r requirements.txt)")
        sys.exit(0)

    tmp = tempfile.mkdtemp(prefix="draft2pptx-gate-")
    md1 = os.path.join(tmp, "TestPlanDraft__doc12__20260904-000000.md")
    with open(md1, "w", encoding="utf-8") as f:
        f.write(DRAFT)

    # ---- CLI: sibling-name default -----------------------------------
    r = subprocess.run(["node", JOB, md1], capture_output=True, text=True, cwd=REPO)
    out1 = md1[:-3] + ".pptx"
    check("converts with the sibling name", r.returncode == 0 and os.path.exists(out1),
          r.stdout + r.stderr)
    check("reports slide count and plan title",
          "slides" in r.stdout and "Route Merge" in r.stdout, r.stdout)

    d = pptx.Presentation(out1)
    slides = list(d.slides)
    texts = [slide_text(s) for s in slides]

    # 1 — the slide walk: title, glance, setup, divider, case, divider,
    #     case, open questions, coverage, issues, closing
    check("11 slides for the sample draft", len(slides) == 11, str(len(slides)))
    check("title slide carries the plan headline",
          "Route Merge" in texts[0] and "TEST PLAN REVIEW" in texts[0], texts[0][:120])
    check("title slide carries the DRAFT pill",
          "DRAFT — MACHINE-GENERATED, UNREVIEWED" in texts[0], texts[0][:200])
    check("title slide carries the Overview facts",
          all(t in texts[0] for t in ("Pro", "3.8", "Claire Wang")), texts[0])
    check("machine banner reaches no slide",
          not any("machine-generated test-plan draft" in t for t in texts), "")

    # 3 — at-a-glance counts (1 positive, 1 negative, 2 real [VERIFY: flags)
    check("glance counts cases and real VERIFY flags",
          "Positive cases" in texts[1] and "Negative cases" in texts[1]
          and "Open [VERIFY] flags" in texts[1], texts[1][:200])
    tiles = [ln for ln in texts[1].split("\n")]
    check("VERIFY tile counts colon-flags only (2, not banner prose)",
          "2" in tiles[tiles.index("Open [VERIFY] flags") - 1]
          if "Open [VERIFY] flags" in tiles else False, str(tiles[:12]))
    check("glance carries the verifier finding",
          "TC-N1 carries a" in texts[1], texts[1])

    # setup checklist
    check("setup items verbatim incl. numbering",
          "1. LRS network with two mergeable routes." in texts[2]
          and "2. Two Pro sessions signed in." in texts[2], texts[2])

    # 4 — case slides (positive at index 4, negative at 6, after dividers)
    check("positive divider names the section",
          "Positive Tests" in texts[3], texts[3])
    case_p = texts[4]
    check("case slide carries id chip, kind chip, title",
          "TC-P1" in case_p and "POSITIVE" in case_p
          and "Merge preserves measures" in case_p, case_p[:200])
    check("case steps verbatim",
          "1. Run Merge Routes on route A and route B." in case_p, case_p)
    check("expected result card",
          "EXPECTED RESULT" in case_p
          and "keeps the source measures unchanged" in case_p, case_p)
    check("trace card",
          "TRACE" in case_p and "the merge must preserve measures" in case_p, case_p)
    check("negative divider carries the CAUTION alert",
          "Negative Tests" in texts[5] and "never the edit succeeding" in texts[5], texts[5])
    check("negative case slide",
          "TC-N1" in texts[6] and "NEGATIVE" in texts[6]
          and "denied with a lock conflict" in texts[6], texts[6][:200])

    # 5 — [VERIFY: ...] spans render amber
    amber = [r.text for s in slides for r in runs_of(s)
             if r.font.color and r.font.color.type is not None
             and str(r.font.color.rgb) == "C2701A" and "[VERIFY:" in r.text]
    check("VERIFY flags surface as amber runs", len(amber) >= 2, str(amber))

    # 6 — native tables with cells intact
    def tables_on(i):
        return [sh.table for sh in slides[i].shapes
                if getattr(sh, "has_table", False) and sh.has_table]

    cov = tables_on(8)
    check("Coverage Map is a native table", len(cov) == 1, str(len(cov)))
    check("Coverage Map cells intact",
          cov and "the merge must preserve measures" in cov[0].rows[1].cells[1].text
          and cov[0].rows[1].cells[2].text.strip() == "TC-P1", "")
    iss = tables_on(9)
    check("Issue Trace is a native table with its row",
          iss and "ArcGISPro/ps-location-referencing#4855" in iss[0].rows[1].cells[0].text,
          "")
    check("Issue Trace addendum note survives (italic aside)",
          "Deterministic addendum" in texts[9], texts[9])

    # 7 — closing provenance
    check("closing slide carries provenance",
          "TestPlanGen prompt v1.7" in texts[10] and "draft2pptx" in texts[10],
          texts[10])

    # 8 — CLI: -o naming, multi-input, usage
    out2 = os.path.join(tmp, "named.pptx")
    r = subprocess.run(["node", JOB, md1, "-o", out2], capture_output=True, text=True, cwd=REPO)
    check("-o names the output", r.returncode == 0 and os.path.exists(out2), r.stdout + r.stderr)
    md2 = os.path.join(tmp, "second.md")
    with open(md2, "w", encoding="utf-8") as f:
        f.write(DRAFT)
    r = subprocess.run(["node", JOB, md1, md2], capture_output=True, text=True, cwd=REPO)
    check("multi-input converts each to a sibling",
          r.returncode == 0 and os.path.exists(md2[:-3] + ".pptx"), r.stdout + r.stderr)
    r = subprocess.run(["node", JOB], capture_output=True, text=True, cwd=REPO)
    check("no-args usage exits nonzero", r.returncode != 0, str(r.returncode))
    r = subprocess.run(["node", JOB, md1, md2, "-o", out2], capture_output=True, text=True, cwd=REPO)
    check("-o with multiple inputs is refused", r.returncode != 0, str(r.returncode))

    # 9 — figures (v1.1): **Figure:** line + --media -> a figure slide
    media = os.path.join(tmp, "media")
    os.makedirs(media, exist_ok=True)
    with open(os.path.join(media, "doc12_slide2_fig1.svg"), "w", encoding="utf-8") as f:
        f.write(FIG_SVG)
    md3 = os.path.join(tmp, "TestPlanDraft__doc12__20260904-000001.md")
    with open(md3, "w", encoding="utf-8") as f:
        f.write(FIG_DRAFT)
    out3 = os.path.join(tmp, "figdeck.pptx")
    r = subprocess.run(["node", JOB, md3, "--media", media, "-o", out3],
                       capture_output=True, text=True, cwd=REPO)
    check("figure draft converts with --media", r.returncode == 0 and os.path.exists(out3),
          r.stdout + r.stderr)
    fd = pptx.Presentation(out3)
    fslides = list(fd.slides)
    ftexts = [slide_text(s) for s in fslides]
    check("figure slide inserted directly after its case (12 slides)",
          len(fslides) == 12 and "TC-P1" in ftexts[4] and "STEPS" in ftexts[4]
          and "FIGURE" in ftexts[5], str(len(fslides)))
    check("figure slide carries id chip, FIGURE tag, and the story's alt text",
          "TC-P1" in ftexts[5] and "FIGURE" in ftexts[5] and FIG_ALT in ftexts[5],
          ftexts[5][:200])
    groups = [sh for sh in fslides[5].shapes
              if sh.shape_type == 6]  # MSO_SHAPE_TYPE.GROUP
    check("cited SVG lands as one native shape group",
          len(groups) == 1 and len(groups[0].shapes) >= 5
          and groups[0].name == "Slide 2 — Routes before the merge",
          f"groups={len(groups)}")
    check("the plate stays dropped and the raw URL reaches no slide",
          all(len(g.shapes) == 6 for g in groups)
          and not any("media/doc12_slide2_fig1.svg" in t for t in ftexts), "")

    # no --media: converts anyway, muted note + stderr coaching
    out4 = os.path.join(tmp, "figdeck-nomedia.pptx")
    r = subprocess.run(["node", JOB, md3, "-o", out4], capture_output=True, text=True, cwd=REPO)
    nd = pptx.Presentation(out4)
    ntexts = [slide_text(s) for s in nd.slides]
    check("without --media the deck still converts, note on the case slide",
          r.returncode == 0 and len(ntexts) == 11
          and f"Figure: {FIG_ALT} (not embedded)" in ntexts[4], r.stderr)
    check("stderr names the --media fix", "--media" in r.stderr, r.stderr)
    # missing file under --media: same degrade, different reason
    empty = os.path.join(tmp, "empty-media")
    os.makedirs(empty, exist_ok=True)
    r = subprocess.run(["node", JOB, md3, "--media", empty, "-o", out4],
                       capture_output=True, text=True, cwd=REPO)
    check("missing file under --media degrades with its own reason",
          r.returncode == 0 and "no such file under --media" in r.stderr, r.stderr)
    # a --media path that is not a directory is a hard usage error
    r = subprocess.run(["node", JOB, md3, "--media", os.path.join(tmp, "nope")],
                       capture_output=True, text=True, cwd=REPO)
    check("--media non-directory refused", r.returncode != 0 and "not a directory" in r.stderr,
          r.stderr)

    print(f"\n{len(PASS)}/{len(PASS) + len(FAIL)} checks passed")
    sys.exit(1 if FAIL else 0)


if __name__ == "__main__":
    main()
