"""Gate for svg2pptx v1.2 — SlideFigures SVG figures -> editable PowerPoint shapes.

Asserts the contract review decks depend on, over a fixture SVG that
exercises the FULL vocabulary SlideFigures emits (every element kind,
both arrow markers, dash rhythms, opacity, rotation, the translate
group, node labels, legend, escaped text):

  1. the output is a well-formed pptx package: every part parses as XML,
     the content types / rels wiring is complete, one slide per SVG;
  2. every SVG primitive lands as a native shape — lines, ticks, extents,
     split hairline + dot, nodes (box/ellipse/diamond), freeform paths,
     text boxes — grouped per figure, nothing rasterised;
  3. styles resolve from the figure's own <style> block: palette hexes on
     strokes/fills, butt vs round caps, dash patterns, the split's
     opacity as stroke alpha, font sizes/weights on runs;
  4. arrowheads become triangle tailEnds (both marker colours);
  5. node labels live INSIDE their shape's text body (editable in place),
     never as loose text boxes;
  6. the figure's <title>/<desc> ride the group as name + alt text;
  7. python-pptx (already in review/harness/requirements.txt) can open
     the package — the closest scriptable stand-in for PowerPoint itself;
  8. v1.2 slide dress: the SVG's plate (white card + border) is DROPPED —
     the slide is the background — and every slide carries a title band:
     the figure's <title> as the slide title plus the source document's
     title above it, looked up from the corpus naming (doc{N}_*.svg ->
     sibling {slug}__doc{N}.md sidecar H1) or forced via --doc-title.

Pure stdlib except that last python-pptx leg, which degrades to a note
when the library is absent. Exit nonzero on any failure.
"""
import os
import re
import subprocess
import sys
import xml.etree.ElementTree as ET
import zipfile

HERE = os.path.dirname(os.path.abspath(__file__))
SCRIPT = os.path.join(HERE, '..', 'svg2pptx.mjs')
failures = []


def check(cond, label):
    print(('ok   ' if cond else 'FAIL ') + label)
    if not cond:
        failures.append(label)


# ---- fixture: the full SlideFigures vocabulary in one figure -------------
# the <style> block is the real figStyle() emission (scripts/SlideFigures.ts)
STYLE = (
    '<style>'
    '.plate{fill:#FFFFFF;stroke:#D7DFDF;stroke-width:1}'
    '.ln{fill:none;stroke-linecap:round;stroke-linejoin:round}'
    '.route{stroke:#16302F;stroke-width:3;stroke-dasharray:10 6;stroke-linecap:butt}'
    '.ctx{stroke:#B9C6C6;stroke-width:2.4}'
    '.event{stroke-width:8}.flat{stroke-linecap:butt}'
    '.tick{stroke:#6E8285;stroke-width:1.15}'
    '.maj{stroke:#4E6265;stroke-width:1.4}'
    '.leader{stroke:#6E8285;stroke-width:1}'
    '.split{stroke:#16302F;stroke-width:1.4;stroke-dasharray:3 2.5;opacity:.55}'
    '.splitdot{fill:#FFFFFF;stroke:#16302F;stroke-width:1.6}'
    '.edge{stroke:#4E6265;stroke-width:1.8}'
    '.free{stroke-width:2.2}'
    '.freefill{stroke-width:2.2;stroke-linejoin:round}'
    '.dashed{stroke-dasharray:7 4.5}.dotted{stroke-dasharray:1.6 3.6}'
    '.node{fill:#FFFFFF;stroke:#16302F;stroke-width:1.6}'
    '.t-plain{fill:#FFFFFF}.t-ink{fill:#E9EDED}.t-muted{fill:#EFF2F2}'
    '.t-cool{fill:#E5F0F5}.t-warm{fill:#F9F0E2}.t-green{fill:#E6F2EC}'
    '.t-violet{fill:#EFEAF7}.t-red{fill:#F8E9E5}'
    '.nlabel{font-size:12px;fill:#16302F;font-weight:500}'
    '.swatch{stroke-width:5}'
    '.legend{font-size:10.5px;fill:#4E6265}'
    '.s-ink{stroke:#16302F}.f-ink{fill:#16302F}'
    '.s-muted{stroke:#6E8285}.f-muted{fill:#6E8285}'
    '.s-cool{stroke:#1B6E8C}.f-cool{fill:#1B6E8C}'
    '.s-warm{stroke:#C2701A}.f-warm{fill:#9C5A12}'
    '.s-green{stroke:#2E7D5B}.f-green{fill:#2E7D5B}'
    '.s-violet{stroke:#7A5AA6}.f-violet{fill:#7A5AA6}'
    '.s-red{stroke:#B2442F}.f-red{fill:#B2442F}'
    '.event.s-cool,.swatch.s-cool{stroke:#4FA7D5}'
    '.event.s-warm,.swatch.s-warm{stroke:#E39A45}'
    '.event.s-green,.swatch.s-green{stroke:#4EB183}'
    '.event.s-violet,.swatch.s-violet{stroke:#A58BD3}'
    '.event.s-red,.swatch.s-red{stroke:#DC8168}'
    "text{font-family:'Segoe UI',system-ui,Roboto,'Helvetica Neue',Arial,sans-serif}"
    '.measure{font-size:11px;fill:#6E8285;font-variant-numeric:tabular-nums}'
    '.id{font-size:12.5px;font-weight:600}.note{font-size:12px;fill:#16302F}'
    '</style>'
    '<defs><marker id="ar" viewBox="0 0 8 8" refX="6" refY="4" markerWidth="4.4" '
    'markerHeight="4.4" orient="auto-start-reverse">'
    '<path d="M0 0.7 L8 4 L0 7.3 z" fill="#16302F"/></marker>'
    '<marker id="ae" viewBox="0 0 8 8" refX="6" refY="4" markerWidth="4.4" '
    'markerHeight="4.4" orient="auto-start-reverse">'
    '<path d="M0 0.7 L8 4 L0 7.3 z" fill="#4E6265"/></marker></defs>'
)

FIXTURE = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 760 320" width="760" '
    'height="320" role="img" aria-label="Slide 3 diagram">'
    '<title>Slide 3 — Route split (1 of 2)</title>'
    '<desc>Schematic of the case: loop route, split at measure 20.</desc>'
    + STYLE +
    '<rect class="plate" x="1" y="1" width="758" height="318" rx="6"/>'
    '<g transform="translate(20,40)">'
    '<line class="ln event flat s-cool" x1="60" y1="60" x2="300" y2="60"/>'
    '<line class="ln event flat s-warm" x1="300" y1="60" x2="540" y2="60"/>'
    '<line class="ln tick maj" x1="60" y1="52.5" x2="60" y2="67.5"/>'
    '<line class="ln tick" x1="120" y1="55.5" x2="120" y2="64.5"/>'
    '<line class="ln route" x1="40" y1="60" x2="640" y2="60"/>'
    '<line class="ln route" x1="632" y1="60" x2="640" y2="60" marker-end="url(#ar)"/>'
    '<line class="split" x1="300" y1="49.5" x2="300" y2="70.5"/>'
    '<circle class="splitdot" cx="300" cy="60" r="3.2"/>'
    '<text class="measure" x="60" y="44.5" text-anchor="middle" '
    'dominant-baseline="central">0</text>'
    '<text class="id f-cool" x="180" y="80" text-anchor="middle" '
    'dominant-baseline="central">E7 &amp; E8</text>'
    '<rect class="node t-cool s-ink" x="80" y="120" width="120" height="48" rx="7"/>'
    '<text class="nlabel" x="140" y="136.5" text-anchor="middle" '
    'dominant-baseline="central">Create route</text>'
    '<text class="nlabel" x="140" y="151.5" text-anchor="middle" '
    'dominant-baseline="central">segment</text>'
    '<ellipse class="node t-warm s-warm" cx="320" cy="144" rx="55" ry="24"/>'
    '<text class="nlabel" x="320" y="144" text-anchor="middle" '
    'dominant-baseline="central">Calibrate</text>'
    '<polygon class="node t-plain s-ink" points="480,120 540,144 480,168 420,144"/>'
    '<line class="ln edge" x1="200" y1="144" x2="265" y2="144" marker-end="url(#ae)"/>'
    '<path class="ln edge dashed" d="M 375 144 L 398 144 L 398 144 L 420 144"/>'
    '<path class="ln free s-violet dotted" '
    'd="M 560 120 C 580 100 600 180 620 140 Q 630 130 640 140"/>'
    '<path class="freefill t-green s-green" d="M 560 220 L 600 200 L 640 220 Z"/>'
    '<line class="ln swatch flat s-cool" x1="40" y1="260" x2="62" y2="260"/>'
    '<text class="legend" x="68" y="260" dominant-baseline="central">'
    'E7 0 → 20</text>'
    '<rect class="node t-violet s-violet" x="560" y="20" width="40" height="80" '
    'rx="7" transform="rotate(9 580 60)"/>'
    '<text class="nlabel" x="580" y="60" text-anchor="middle" '
    'dominant-baseline="central">Publish</text>'
    '</g></svg>'
)

open('fix_slide3.svg', 'w', encoding='utf-8').write(FIXTURE)
open('fix_slide10.svg', 'w', encoding='utf-8').write(
    FIXTURE.replace('Slide 3', 'Slide 10'))

out = subprocess.run(
    ['node', SCRIPT, 'fix_slide3.svg', 'fix_slide10.svg', '-o', 'fix_figures.pptx'],
    capture_output=True, text=True, encoding='utf-8')
check(out.returncode == 0, 'converter exits 0 (' + out.stderr.strip()[:300] + ')')
check('2 figure(s)' in out.stdout, 'summary line reports both figures')
if out.returncode != 0:
    print('RESULT: FAIL')
    sys.exit(1)

z = zipfile.ZipFile('fix_figures.pptx')
names = set(z.namelist())

# ---- 1: package shape ----------------------------------------------------
need = ['[Content_Types].xml', '_rels/.rels', 'ppt/presentation.xml',
        'ppt/_rels/presentation.xml.rels', 'ppt/slideMasters/slideMaster1.xml',
        'ppt/slideLayouts/slideLayout1.xml', 'ppt/theme/theme1.xml',
        'ppt/slides/slide1.xml', 'ppt/slides/slide2.xml',
        'ppt/slides/_rels/slide1.xml.rels', 'docProps/core.xml']
for p in need:
    check(p in names, f'part present: {p}')
for p in sorted(names):
    try:
        ET.fromstring(z.read(p))
        ok = True
    except ET.ParseError:
        ok = False
    check(ok, f'{p}: parses as XML')
ct = z.read('[Content_Types].xml').decode('utf-8')
check('slide+xml' in ct and ct.count('slides/slide') == 2,
      'content types declare both slides')
pres = z.read('ppt/presentation.xml').decode('utf-8')
check(pres.count('<p:sldId ') == 2, 'presentation lists two slides')
check('cx="12192000" cy="6858000"' in pres, 'slide size is 16:9')

s1 = z.read('ppt/slides/slide1.xml').decode('utf-8')
s2 = z.read('ppt/slides/slide2.xml').decode('utf-8')

# ---- 2: every primitive becomes a native shape ---------------------------
check(s1.count('<p:grpSp>') == 1, 'one group per figure')
check('prst="line"' in s1, 'lines land as line shapes')
check('prst="roundRect"' in s1, 'box node lands as roundRect')
check('prst="ellipse"' in s1, 'ellipse node + splitdot land as ellipses')
check('prst="diamond"' in s1, '4-point node polygon lands as a diamond')
check('<a:custGeom>' in s1, 'freeform paths land as custGeom')
check('<a:cubicBezTo>' in s1, 'curves survive as cubic beziers')
check(s1.count('<a:close/>') >= 1, 'closed freeform stays closed')
check('name="measure"' in s1 and 'name="tick' not in s1.replace('name="tick"', '')
      or 'name="tick"' in s1, 'shapes named by role for the selection pane')
check('<a:blip' not in s1 and 'image' not in ct, 'nothing rasterised')

# ---- 8 (v1.2): no plate — the slide is the background --------------------
check('name="plate"' not in s1, 'plate: not emitted as a shape')
check('D7DFDF' not in s1, 'plate: its border colour appears nowhere')
check(s1.count('prst="roundRect"') == 2,
      'plate: only the two box nodes round-rect (no background card)')

# ---- 3: styles resolve from the figure's own stylesheet ------------------
check('val="16302F"' in s1, 'route stroke: ink resolved')
check('val="4FA7D5"' in s1 and 'val="E39A45"' in s1,
      'extent strokes: SOFT band variants resolved via compound rules (v1.9)')
check('val="7A5AA6"' in s1,
      'freeform keeps the DEEP violet — compound rules scope to event/swatch only')
check('val="E5F0F5"' in s1 and 'val="F9F0E2"' in s1,
      'node fills: palette tints resolved')
m = re.search(r'<a:ln w="(\d+)" cap="flat">'
              r'<a:solidFill><a:srgbClr val="4FA7D5"/>', s1)
check(m is not None, 'event extent: butt cap carried through')
if m:
    check(abs(int(m.group(1)) - round(8 * 9525)) <= 10,
          f'event extent: 8px band -> EMU width ({m.group(1)})')
check(re.search(r'val="16302F"><a:alpha val="55000"/>', s1) is not None,
      "split hairline: opacity .55 -> 55% stroke alpha")
check('<a:prstDash val="sysDash"/>' in s1, 'split dash (3 2.5) -> sysDash')
check(re.search(r'<a:ln w="\d+" cap="flat"><a:solidFill><a:srgbClr val="16302F"/>'
                r'</a:solidFill><a:prstDash val="dash"/>', s1) is not None,
      'route: dash rhythm + butt cap carried onto the shape (v1.7 style)')
check('<a:prstDash val="dash"/>' in s1, 'dashed (7 4.5) -> dash')
check('<a:prstDash val="sysDot"/>' in s1, 'dotted (1.6 3.6) -> sysDot')
check(re.search(r'rot="540000"[^>]*>', s1) is not None,
      'rotate(9 ...) -> native 9-degree shape rotation')

# ---- 4: arrowheads -------------------------------------------------------
check(s1.count('<a:tailEnd type="triangle"') == 2,
      'both marker-ends -> triangle tailEnds (route + edge)')

# ---- 5: node labels live inside their shapes -----------------------------
mnode = re.search(r'<p:sp>(?:(?!</p:sp>).)*roundRect(?:(?!</p:sp>).)*Create route'
                  r'(?:(?!</p:sp>).)*</p:sp>', s1, re.S)
check(mnode is not None, 'box node carries its label in its own text body')
if mnode:
    check('segment' in mnode.group(0), 'two-row label -> two paragraphs, one shape')
check('Publish' in s1 and 'Calibrate' in s1, 'every node label carried through')
check(s1.count('Create route') == 1, 'node label not duplicated as a text box')
check('wrap="none"' in s1, 'loose text boxes stay one line (wrap off)')
check('E7 &amp; E8' in s1, 'escaped text content survives escaped')
check('sz="825"' in s1, 'measure text: 11px -> 8.25pt run size')
check(re.search(r'sz="9\d\d" b="1"', s1) is not None,
      'event id: 12.5px semibold -> bold run')
check('typeface="Segoe UI"' in s1, 'font family read from the stylesheet')

# ---- 6: title/desc ride the group ----------------------------------------
check('Slide 3 — Route split (1 of 2)' in s1, 'group named from <title>')
check('Schematic of the case' in s1, '<desc> lands as group alt text')
check('Slide 10' in s2, 'second SVG landed on the second slide')

# ---- 8 (v1.2): title band -------------------------------------------------
check('<a:t>Slide 3 — Route split (1 of 2)</a:t>' in s1,
      'slide title: the figure <title> lands as a real title text box')
check('name="slide title"' in s1, 'slide title box named for the selection pane')
check('name="document title"' not in s1,
      'no doc{N}_ prefix and no --doc-title -> no document-title box')
gym = re.search(r'name="Slide 3[^"]*"[^>]*/>.*?<a:off x="\d+" y="(\d+)"/>', s1, re.S)
check(gym is not None and int(gym.group(1)) >= 1120140,
      'figure group sits below the title band')

# doc-title lookup: media naming doc{N}_*.svg -> sibling kind folder's
# {slug}__doc{N}.md sidecar, H1 wins
os.makedirs('fixlib/media', exist_ok=True)
os.makedirs('fixlib/Test Plans', exist_ok=True)
open('fixlib/media/doc7_slide3.svg', 'w', encoding='utf-8').write(FIXTURE)
open('fixlib/Test Plans/route-split-cases__doc7.md', 'w', encoding='utf-8').write(
    '# Route Split Cases\n\nbody\n')
out2 = subprocess.run(
    ['node', SCRIPT, 'fixlib/media', '-o', 'fix_doc.pptx'],
    capture_output=True, text=True, encoding='utf-8')
check(out2.returncode == 0, 'doc-title leg: converter exits 0')
zd = zipfile.ZipFile('fix_doc.pptx')
sd = zd.read('ppt/slides/slide1.xml').decode('utf-8')
check('<a:t>Route Split Cases</a:t>' in sd,
      'document title: sidecar H1 found via the doc7_ prefix')
check('name="document title"' in sd, 'document title box named')
out3 = subprocess.run(
    ['node', SCRIPT, 'fixlib/media', '--doc-title', 'Override Deck',
     '-o', 'fix_doc2.pptx'],
    capture_output=True, text=True, encoding='utf-8')
zo = zipfile.ZipFile('fix_doc2.pptx')
so = zo.read('ppt/slides/slide1.xml').decode('utf-8')
check(out3.returncode == 0 and '<a:t>Override Deck</a:t>' in so
      and 'Route Split Cases' not in so,
      '--doc-title overrides the sidecar lookup')

# ---- 7: python-pptx opens it (PowerPoint's scriptable stand-in) ----------
try:
    from pptx import Presentation
    prs = Presentation('fix_figures.pptx')
    check(len(prs.slides) == 2, 'python-pptx: opens, two slides')
    shp = prs.slides[0].shapes
    check(len(shp) == 2 and shp[0].shape_type is not None,
          'python-pptx: slide holds the title box and the figure group')
    grp = shp[-1]
    check(grp.name.startswith('Slide 3'), 'python-pptx: group name readable')
    check(len(grp.shapes) >= 20,
          f'python-pptx: group holds the shapes ({len(grp.shapes)})')
except ImportError:
    print('note python-pptx not installed - open-leg skipped')

print()
if failures:
    print(f'RESULT: FAIL - {len(failures)} assertion(s) failed')
    sys.exit(1)
print('RESULT: PASS')
