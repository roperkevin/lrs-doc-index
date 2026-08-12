"""Gate for pptxgen/ — the sidecar-markdown -> .pptx briefing-deck converter.

Runs the real CLI (pptxgen/sidecar_to_pptx.py) over the sample sidecars
rendered by render_sample.py and asserts the resulting decks' structure
by reading them back with python-pptx. Everything runs offline: the
enhanced rendering path is exercised from the checked-in
enhanced_outline_fixture.json, and the --enhance stage is tested via its
REQUIRED failure behavior (no API key -> deterministic fallback with a
warning; --strict-enhance -> exit 3). No network, no anthropic install.

Suites:
  1  deterministic conversion of sample_sidecar.md / _related.md:
     16:9 geometry, verbatim title + doc_kind chip, summary + one chip
     per keyword, every parsed section title exactly once as a base
     content title in parse order (extras only "(cont.)"-suffixed),
     speaker notes round-trip off-canvas, GFM table -> native table with
     matching column count, planted image -> Picture shape, missing
     images -> styled placeholder, related closer only when the sidecar
     has related entries (all titles + why strings).
  2  outline contract: outline_from_sidecar() round-trips through the
     OUTLINE_SCHEMA field subset and load_outline() (keeps the mapping,
     the schema, and the validator in lockstep).
  3  enhanced path (offline): enhanced_outline_fixture.json validates
     and renders via --from-outline with agenda / persona / takeaways.
  4  docx-lane tolerance: "##"-heading bodies convert without slide
     kickers.
  5  overflow: 40 bullets + a 30-row table paginate into "(cont.)"
     slides; no rendered table exceeds 1+8 rows; header repeats.
  6  degenerate inputs: header-only sidecar renders title+summary;
     non-sidecar input exits 2.
  7  --enhance fallback semantics (above).

Prereqs: render_sample.py has run (sample_sidecar.md / _related.md).
This gate also (re)plants review/media/doc42_image1.png so exactly one
of the sample's image links resolves — the rest exercise placeholders.
"""
import json
import os
import re
import subprocess
import sys

sys.path.insert(0, os.path.join('..', '..', 'pptxgen'))
from deck_outline import OUTLINE_SCHEMA, load_outline, outline_from_sidecar
from sidecar_parser import parse_sidecar

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

CLI = os.path.join('..', '..', 'pptxgen', 'sidecar_to_pptx.py')
ok = True


def check(cond, label):
    global ok
    print(('ok   ' if cond else 'FAIL ') + label)
    if not cond:
        ok = False
    return cond


def run_cli(args, env=None):
    return subprocess.run([sys.executable, CLI, *args],
                          capture_output=True, text=True, env=env)


def slide_texts(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if t.strip():
                out.append(t)
        if shape.has_table:
            for row in shape.table.rows:
                out.append('|'.join(c.text for c in row.cells))
    return out


def slide_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.split('\n')[0].strip()
    return ''


def tables(prs):
    return [s.table for sl in prs.slides for s in sl.shapes if s.has_table]


def pictures(prs):
    return [s for sl in prs.slides for s in sl.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


# --------------------------------------------------------------- prereqs
for f in ('sample_sidecar.md', 'sample_sidecar_related.md'):
    if not os.path.exists(f):
        print(f'FAIL missing {f} — run render_sample.py first')
        sys.exit(1)

# plant one real image (inline clone of make_fixtures.make_png — importing
# make_fixtures would re-run the whole fixture build as a side effect)
import random  # noqa: E402
import struct  # noqa: E402
import zlib  # noqa: E402


def make_png(w=40, h=40):
    def chunk(tag, data):
        c = tag + data
        return struct.pack('>I', len(data)) + c \
            + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
    raw = b''.join(b'\x00' + bytes(random.randrange(256) for _ in range(w * 3))
                   for _ in range(h))
    return (b'\x89PNG\r\n\x1a\n'
            + chunk(b'IHDR', struct.pack('>IIBBBBB', w, h, 8, 2, 0, 0, 0))
            + chunk(b'IDAT', zlib.compress(raw))
            + chunk(b'IEND', b''))


os.makedirs(os.path.join('..', 'media'), exist_ok=True)
with open(os.path.join('..', 'media', 'doc42_image1.png'), 'wb') as fh:
    fh.write(make_png(120, 60))

md = open('sample_sidecar_related.md', encoding='utf-8').read()
sc = parse_sidecar(md)

# ---------------------------------------------- 1. deterministic decks
for name, expect_related in (('sample_sidecar.md', False),
                             ('sample_sidecar_related.md', True)):
    out = os.path.splitext(name)[0].replace('sample_sidecar',
                                            'sample_deck') + '.pptx'
    r = run_cli([name, '-o', out, '--media-dir', '../media'])
    check(r.returncode == 0, f'{name}: CLI exit 0')
    prs = Presentation(out)
    this_sc = parse_sidecar(open(name, encoding='utf-8').read())

    check(prs.slide_width == 12192000 and prs.slide_height == 6858000,
          f'{out}: 16:9 slide size')

    s1 = prs.slides[0]
    check(any(this_sc.title in t for t in slide_texts(s1)),
          f'{out}: slide 1 carries the yaml title verbatim')
    check(any(t.strip() == str(this_sc.meta['doc_kind']).upper()
              for t in slide_texts(s1)),
          f'{out}: slide 1 carries the doc_kind chip')

    s2 = prs.slides[1]
    check(slide_title(s2) == 'Summary' and
          any(this_sc.summary[:60] in t for t in slide_texts(s2)),
          f'{out}: slide 2 is the summary slide')
    kw_ok = all(any(t.strip() == kw for t in slide_texts(s2))
                for kw in this_sc.meta['keywords'])
    check(kw_ok, f'{out}: one chip per keyword on the summary slide')

    # section titles appear exactly once as base titles, in parse order,
    # with continuations only ever "(cont.)"-suffixed
    content = list(prs.slides)[2:]
    if expect_related:
        rel_start = next(i for i, sl in enumerate(content)
                         if slide_title(sl).startswith('Related documents'))
        related_slides, content = content[rel_start:], content[:rel_start]
    titles = [slide_title(sl) for sl in content]
    expected = [s.title for s in this_sc.sections]
    base = [t for t in titles if not t.endswith('(cont.)')]
    check(base == expected,
          f'{out}: section titles once each, in parse order '
          f'({len(base)}/{len(expected)})')
    conts = [t for t in titles if t.endswith('(cont.)')]
    check(all(t[:-len(' (cont.)')] in expected for t in conts),
          f'{out}: extra content slides only ever "(cont.)"-suffixed '
          f'({len(conts)} continuations)')

    # notes round-trip: on the section's first slide, never on canvas
    notes_ok = canvas_ok = True
    for section in this_sc.sections:
        if not section.notes:
            continue
        idx = titles.index(section.title)
        sl = content[idx]
        note_text = sl.notes_slide.notes_text_frame.text \
            if sl.has_notes_slide else ''
        if section.notes[0] not in note_text:
            notes_ok = False
        if any(section.notes[0] in t for t in slide_texts(sl)):
            canvas_ok = False
    check(notes_ok, f'{out}: "### Notes" text lands in speaker notes')
    check(canvas_ok, f'{out}: notes text never appears on the canvas')

    tbls = tables(prs)
    first_table = next(b for s in this_sc.sections for b in s.blocks
                       if type(b).__name__ == 'Table')
    check(any(len(t.columns) == len(first_table.rows[0]) for t in tbls),
          f'{out}: native table with the md table column count')
    check(all(len(t.rows) <= 9 for t in tbls),
          f'{out}: no rendered table exceeds header + 8 body rows')

    check(len(pictures(prs)) >= 1,
          f'{out}: planted PNG embedded as a Picture shape')
    all_text = [t for sl in prs.slides for t in slide_texts(sl)]
    check(any('Image: doc42_' in t for t in all_text),
          f'{out}: missing images render the styled placeholder')

    if expect_related:
        rel_text = [t for sl in related_slides for t in slide_texts(sl)]
        got = all(any(e.title in t for t in rel_text) and
                  any(e.why[:40] in t for t in rel_text)
                  for e in this_sc.related)
        check(got, f'{out}: related closer lists all titles + why strings')
    else:
        check(not any(slide_title(sl).startswith('Related documents')
                      for sl in prs.slides),
              f'{out}: no related slide when the list is empty')

# ------------------------------------- 2. outline contract round-trip
det = outline_from_sidecar(sc)
subset = {k: det.to_dict()[k] for k in OUTLINE_SCHEMA['required']}
subset = json.loads(json.dumps(subset))
try:
    det2 = load_outline(subset, sc)
    rt = [s.kind for s in det2.slides] == [s.kind for s in det.slides] \
        and det2.related == det.related
except Exception as exc:  # noqa: BLE001
    rt = False
    print('     ', exc)
check(rt, 'outline contract: deterministic mapping round-trips through '
          'the schema subset + load_outline')

# ------------------------------------------- 3. enhanced path (offline)
fixture = json.load(open('enhanced_outline_fixture.json', encoding='utf-8'))
fixture = {k: v for k, v in fixture.items() if not k.startswith('_')}
try:
    enh = load_outline(fixture, sc)
    check(True, 'enhanced fixture validates via load_outline')
except Exception as exc:  # noqa: BLE001
    check(False, f'enhanced fixture validates via load_outline ({exc})')

with open('enhanced_outline_clean.json', 'w', encoding='utf-8') as fh:
    json.dump(fixture, fh)
r = run_cli(['sample_sidecar_related.md', '-o', 'sample_deck_enhanced.pptx',
             '--media-dir', '../media',
             '--outline-json', 'enhanced_outline_clean.json',
             '--from-outline'])
check(r.returncode == 0, 'enhanced fixture renders via --from-outline')
prs = Presentation('sample_deck_enhanced.pptx')
all_text = [t for sl in prs.slides for t in slide_texts(sl)]
kinds = [s.kind for s in enh.slides]
check('agenda' in kinds and
      any('Why locks must move to route creation' in t for t in all_text),
      'enhanced deck: agenda slide with fixture items')
persona = next(b.text for s in enh.slides for b in s.blocks
               if b.kind == 'persona')
check(any(persona in t for t in all_text),
      'enhanced deck: persona callout text on canvas')
take = next(s for s in enh.slides if s.kind == 'takeaways')
check(all(any(item in t for t in all_text) for item in take.items),
      'enhanced deck: all takeaways rendered')
check(any(slide_title(sl).startswith('Related documents')
          for sl in prs.slides),
      'enhanced deck: related closer present')

# ------------------------------------------ 4. docx-lane tolerance
header = md[:md.index('\n---\n', md.index('## Summary')) + len('\n---\n')]
docx_body = """
## Requirements overview

The route lock must be acquired at creation time.

### Locking rules

- acquire on create
  - covers Extend and Realign
- release on save

## Open questions

| Question | Owner |
| --- | --- |
| Lock timeout | PE |
"""
open('docx_body_sidecar.md', 'w', encoding='utf-8').write(header + docx_body)
r = run_cli(['docx_body_sidecar.md', '-o', 'docx_body_deck.pptx'])
prs = Presentation('docx_body_deck.pptx')
titles = [slide_title(sl) for sl in prs.slides]
check(r.returncode == 0 and 'Requirements overview' in titles and
      'Open questions' in titles,
      'docx lane: one content slide per "##" heading')
check(not any(re.fullmatch(r'SLIDE \d+', t)
              for sl in prs.slides for t in slide_texts(sl)),
      'docx lane: no SLIDE N kickers')
check(any('Locking rules' in t
          for sl in prs.slides for t in slide_texts(sl)),
      'docx lane: deeper headings render as sub-headers')

# ------------------------------------------------- 5. overflow
bullets = '\n'.join(f'  - overflow bullet {i} with some padding text words'
                    for i in range(40))
rows = '\n'.join(f'| r{i}c0 | r{i}c1 | r{i}c2 |' for i in range(30))
overflow_body = f"""
## Slide 1 — Bullet storm

{bullets}

## Slide 2 — Table storm

| h0 | h1 | h2 |
| --- | --- | --- |
{rows}
"""
open('overflow_sidecar.md', 'w', encoding='utf-8').write(header + overflow_body)
r = run_cli(['overflow_sidecar.md', '-o', 'overflow_deck.pptx'])
prs = Presentation('overflow_deck.pptx')
titles = [slide_title(sl) for sl in prs.slides]
check(r.returncode == 0 and
      titles.count('Bullet storm (cont.)') >= 1 and
      titles.count('Table storm (cont.)') >= 2,
      'overflow: long sections continue on "(cont.)" slides')
tbls = tables(prs)
check(tbls and all(len(t.rows) <= 9 for t in tbls),
      'overflow: every table part capped at header + 8 body rows')
check(all(t.cell(0, 0).text == 'h0' for t in tbls),
      'overflow: header row repeated on every table part')

# ------------------------------------------------ 6. degenerate inputs
open('headeronly_sidecar.md', 'w', encoding='utf-8').write(header + '\n')
r = run_cli(['headeronly_sidecar.md', '-o', 'headeronly_deck.pptx'])
prs = Presentation('headeronly_deck.pptx')
check(r.returncode == 0 and len(prs.slides._sldIdLst) >= 2 and
      slide_title(prs.slides[1]) == 'Summary',
      'degenerate: header-only sidecar renders title + summary, exit 0')

open('notasidecar.md', 'w', encoding='utf-8').write('# just markdown\n')
r = run_cli(['notasidecar.md', '-o', 'nope.pptx'])
check(r.returncode == 2 and not os.path.exists('nope.pptx'),
      'degenerate: non-sidecar input exits 2, writes nothing')

# ------------------------------------------- 7. --enhance fallback
env = {k: v for k, v in os.environ.items() if k != 'ANTHROPIC_API_KEY'}
r = run_cli(['sample_sidecar.md', '-o', 'fallback_deck.pptx', '--enhance'],
            env=env)
check(r.returncode == 0 and 'falling back' in r.stderr and
      os.path.exists('fallback_deck.pptx'),
      '--enhance without a key: warns and falls back, deck still built')
r = run_cli(['sample_sidecar.md', '-o', 'strict_deck.pptx', '--enhance',
             '--strict-enhance'], env=env)
check(r.returncode == 3 and not os.path.exists('strict_deck.pptx'),
      '--enhance --strict-enhance without a key: exit 3, no deck')

print()
print('RESULT:', 'PASS' if ok else 'FAIL')
sys.exit(0 if ok else 1)
