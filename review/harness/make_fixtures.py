"""Build OOXML fixtures for the harness.

Produces real Office files via python-pptx / python-docx (independent zip
writer, genuine OOXML), records the planted ground-truth tokens for the
recall check, and emits .b64 companions for the Node runners.

Since the v1.7 formatting work the fixtures also plant *structure* —
slide titles, outline-leveled paragraphs, docx headings and numPr list
items, and a sheets.json workbook stand-in — recorded in
planted_format.json for check_format.py. run_diff.py (the v1.5/v1.6
equivalence gate) is unaffected: both versions see the same fixtures.

Since the v1.8 core-properties work the two "real" fixtures also plant
authorship core properties (author with an ampersand, last-modified-by
with a smart apostrophe — both exercise entity decoding — and a fixed
modified timestamp), recorded under each fixture's 'core' key in
planted_format.json, plus a noprops_deck.pptx.b64 companion (edge_deck
with docProps/core.xml stripped) for the degrades-to-empty case.
"""
import base64
import datetime
import io
import json
import os
import random
import struct
import zipfile
import zlib

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches

random.seed(99)

VOCAB = ['route', 'measure', 'referent', 'calibration point', 'centerline', 'event',
         'LRS Network', 'segment', 'Merge Centerlines', 'Retire Routes', 'geoprocessing',
         'Straight Line Diagram', 'vertex spacing', 'positive case', 'negative case',
         'expected result', 'applyEdits', 'dominance', 'concurrency', 'Stay Put', 'Move',
         'Retire', 'Snap', 'Cover', 'iteration', 'estimate', 'assignment']


def sent(n):
    return ' '.join(random.choice(VOCAB) for _ in range(n))


def make_png(w=40, h=40):
    def chunk(tag, data):
        c = tag + data
        return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
    raw = b''.join(b'\x00' + bytes(random.randrange(256) for _ in range(w * 3)) for _ in range(h))
    return (b'\x89PNG\r\n\x1a\n'
            + chunk(b'IHDR', struct.pack('>IIBBBBB', w, h, 8, 2, 0, 0, 0))
            + chunk(b'IDAT', zlib.compress(raw))
            + chunk(b'IEND', b''))


tokens = {}
fmt = {}

# fixture 1: realistic test-plan deck (tables, notes, images, issue urls)
prs = Presentation()
planted = []
deck_fmt = {'titles': [], 'notes_count': 0, 'lvl1': [], 'lvl2': []}
for i in range(18):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    title = f"Slide{i} " + sent(4)
    s.shapes.title.text = title
    planted += title.split()
    deck_fmt['titles'].append(title)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2)).text_frame
    for _ in range(4):
        p = tb.add_paragraph()
        p.text = sent(12)
        planted += p.text.split()
    # planted outline structure: explicit levels -> nested markdown lists
    p = tb.add_paragraph()
    p.text = f"level one item {i} " + sent(3)
    p.level = 1
    planted += p.text.split()
    deck_fmt['lvl1'].append(p.text)
    p = tb.add_paragraph()
    p.text = f"level two item {i} " + sent(2)
    p.level = 2
    planted += p.text.split()
    deck_fmt['lvl2'].append(p.text)
    tbl = s.shapes.add_table(4, 3, Inches(0.5), Inches(4), Inches(8), Inches(2)).table
    for r in range(4):
        for c in range(3):
            v = f"cell{i}r{r}c{c} " + random.choice(VOCAB)
            tbl.cell(r, c).text = v
            planted += v.split()
    notes = s.notes_slide.notes_text_frame
    notes.text = (f"notes{i} PE: Claire Wang Dev: Ito — "
                  f"devtopia.esri.com/ArcGISPro/ps-location-referencing/issues/{4000 + i}")
    planted += notes.text.split()
    deck_fmt['notes_count'] += 1
    if i % 3 == 0:
        s.shapes.add_picture(io.BytesIO(make_png()), Inches(8.6), Inches(0.2), Inches(0.6), Inches(0.6))
CORE_MODIFIED = datetime.datetime(2026, 7, 31, 18, 22, 4)
prs.core_properties.author = 'Claire Wang & Team'
prs.core_properties.last_modified_by = 'Miguel O’Brien'
prs.core_properties.modified = CORE_MODIFIED
prs.save('real_deck.pptx')
tokens['real_deck.pptx'] = planted
deck_fmt['core'] = {'author': 'Claire Wang & Team',
                    'lastEditedBy': 'Miguel O’Brien',
                    'lastEdited': '2026-07-31T18:22:04Z'}
fmt['real_deck.pptx'] = deck_fmt

# fixture 2: docx with headings, nested + merged tables, unicode, issue refs, image
doc = Document()
planted = []
doc_fmt = {'title': None, 'headings': {'1': [], '2': [], '3': []},
           'lists': {'0': [], '1': [], '2': []}}
h = doc.add_heading('Merge Centerlines Test Plan V4 — café naïve 中文 🚀', 0)
planted += h.text.split()
doc_fmt['title'] = h.text


def add_list_item(text, ilvl):
    """Plant a direct-formatting numPr paragraph (how Word marks list
    items on the paragraph itself, unlike python-docx's style-only
    List Bullet)."""
    p = doc.add_paragraph(text)
    ppr = p._p.get_or_add_pPr()
    numpr = ppr.makeelement(qn('w:numPr'), {})
    ilvl_el = ppr.makeelement(qn('w:ilvl'), {qn('w:val'): str(ilvl)})
    numid = ppr.makeelement(qn('w:numId'), {qn('w:val'): '1'})
    numpr.append(ilvl_el)
    numpr.append(numid)
    ppr.append(numpr)
    return p


for i in range(60):
    p = doc.add_paragraph(sent(14))
    planted += p.text.split()
    if i % 20 == 0:
        for lvl in (1, 2, 3):
            hh = doc.add_heading(f"Heading {lvl} section {i} " + sent(2), lvl)
            planted += hh.text.split()
            doc_fmt['headings'][str(lvl)].append(hh.text)
        for ilvl in (0, 1, 2):
            li = add_list_item(f"list item ilvl{ilvl} block {i} " + sent(2), ilvl)
            planted += li.text.split()
            doc_fmt['lists'][str(ilvl)].append(li.text)
t = doc.add_table(rows=5, cols=4)
t.style = 'Table Grid'
for r in range(5):
    for c in range(4):
        v = f"d{r}{c} " + random.choice(VOCAB)
        t.cell(r, c).text = v
        planted += v.split()
t.cell(0, 0).merge(t.cell(0, 1))            # merged cells -> gridSpan
inner = t.cell(2, 2).add_table(rows=2, cols=2)  # nested table
for r in range(2):
    for c in range(2):
        v = f"nest{r}{c}"
        inner.cell(r, c).text = v
        planted += [v]
p = doc.add_paragraph("see #4855 and "
                      "devtopia.esri.com/Beijing-R-D-Center/ExperienceBuilder-Web-Extensions/issues/26161")
planted += p.text.split()
doc.add_picture(io.BytesIO(make_png(60, 60)))
doc.core_properties.author = 'Claire Wang & Team'
doc.core_properties.last_modified_by = 'Miguel O’Brien'
doc.core_properties.modified = CORE_MODIFIED
doc.save('real_doc.docx')
tokens['real_doc.docx'] = planted
doc_fmt['core'] = {'author': 'Claire Wang & Team',
                   'lastEditedBy': 'Miguel O’Brien',
                   'lastEdited': '2026-07-31T18:22:04Z'}
fmt['real_doc.docx'] = doc_fmt

# fixture 3: edge cases — blank slide, notes-only slide
prs2 = Presentation()
prs2.slides.add_slide(prs2.slide_layouts[6])
s2 = prs2.slides.add_slide(prs2.slide_layouts[6])
s2.notes_slide.notes_text_frame.text = "only notes here referent"
prs2.save('edge_deck.pptx')
tokens['edge_deck.pptx'] = ['only', 'notes', 'here', 'referent']
fmt['edge_deck.pptx'] = {'titles': [], 'notes_count': 1, 'lvl1': [], 'lvl2': []}

# fixture 3b: edge_deck with docProps/core.xml stripped — the v1.8
# core-properties degrades-to-empty case (b64 companion only; not in
# the tokens/recall set)
with zipfile.ZipFile('edge_deck.pptx') as zin, \
        zipfile.ZipFile('noprops_deck.pptx', 'w', zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        if item.filename != 'docProps/core.xml':
            zout.writestr(item, zin.read(item.filename))
with open('noprops_deck.pptx', 'rb') as fh:
    open('noprops_deck.pptx.b64', 'w', encoding='utf-8').write(base64.b64encode(fh.read()).decode())

# fixture 4: workbook stand-in for the WorkbookDump mock runner
# (wrap_workbook.py) — a 30-column row exercising the COLCAP cut, a
# pipe-bearing cell exercising escaping, a >300-char cell exercising
# CELLCAP truncation, and an empty sheet.
wide_header = [f"col{c}" for c in range(30)]
wide_row = [f"w{c}" for c in range(30)]
sheets = {
    'Schedule': [wide_header, wide_row, ['w0', 'w1']],
    'Notes': [
        ['Task', 'Detail'],
        ['pipes | in | cells', 'x' * 320],
        ['plain', 'referent calibration'],
    ],
    'Blank': [],
}
json.dump(sheets, open('sheets.json', 'w', encoding='utf-8'))

for f in tokens:
    with open(f, 'rb') as fh:
        open(f + '.b64', 'w', encoding='utf-8').write(base64.b64encode(fh.read()).decode())
json.dump(tokens, open('planted_tokens.json', 'w', encoding='utf-8'))
json.dump(fmt, open('planted_format.json', 'w', encoding='utf-8'))
print({f: os.path.getsize(f) for f in tokens})

# ---- v1.9 batch fixtures (check_batch.py) --------------------------------
# New-behavior fixtures for the SC-2..SC-14/FL-5 script batch. Not part
# of the recall set or the v1.8 check_format list — check_batch.py owns
# their assertions until the batch is promoted.
import re as _re


def _rezip(path, transform):
    """Round-trip a zip, letting transform(name, bytes) edit members."""
    with zipfile.ZipFile(path) as zin:
        items = [(i.filename, zin.read(i.filename)) for i in zin.infolist()]
    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in items:
            zout.writestr(name, transform(name, data))


def _b64(path):
    with open(path, 'rb') as fh:
        open(path + '.b64', 'w', encoding='utf-8').write(base64.b64encode(fh.read()).decode())


# SC-2: reordered_deck — sldIdLst reversed vs part numbering
prs_r = Presentation()
for i in range(4):
    s = prs_r.slides.add_slide(prs_r.slide_layouts[5])
    s.shapes.title.text = f"OrderTitle{i}"
    s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(1)).text_frame.text = f"orderbody{i}"
prs_r.save('reordered_deck.pptx')


def _reverse_sldidlst(name, data):
    if name != 'ppt/presentation.xml':
        return data
    xml = data.decode('utf-8')
    m = _re.search(r'(<p:sldIdLst>)(.*?)(</p:sldIdLst>)', xml, _re.S)
    ids = _re.findall(r'<p:sldId [^>]*/>', m.group(2))
    return (xml[:m.start()] + m.group(1) + ''.join(reversed(ids)) + m.group(3)
            + xml[m.end():]).encode('utf-8')


_rezip('reordered_deck.pptx', _reverse_sldidlst)
_b64('reordered_deck.pptx')

# SC-3: merged_deck — pptx table with a real hMerge continuation cell
prs_m = Presentation()
s_m = prs_m.slides.add_slide(prs_m.slide_layouts[6])
tbl_m = s_m.shapes.add_table(3, 3, Inches(0.5), Inches(1), Inches(8), Inches(2)).table
for r in range(3):
    for c in range(3):
        tbl_m.cell(r, c).text = f"m{r}{c}"
tbl_m.cell(0, 0).text = 'mergedhead'
# spec-shaped horizontal merge: origin carries gridSpan, covered cell
# stays in the markup flagged hMerge (what the SC-3 fix must skip)
tr0 = tbl_m._tbl.tr_lst[0]
tcs = tr0.findall(qn('a:tc'))
tcs[0].set('gridSpan', '2')
tcs[1].set('hMerge', '1')
prs_m.save('merged_deck.pptx')
_b64('merged_deck.pptx')

# SC-4: bigimg_deck — one referenced small png, one referenced over-cap png
prs_b = Presentation()
s_b = prs_b.slides.add_slide(prs_b.slide_layouts[6])
s_b.shapes.add_picture(io.BytesIO(make_png()), Inches(0.5), Inches(0.5), Inches(1), Inches(1))
s_b.shapes.add_picture(io.BytesIO(make_png(450, 300)), Inches(2), Inches(0.5), Inches(3), Inches(2))
prs_b.save('bigimg_deck.pptx')
_b64('bigimg_deck.pptx')

# SC-6: relswap_deck — slide rels rewritten with Target BEFORE Id
prs_s = Presentation()
s_s = prs_s.slides.add_slide(prs_s.slide_layouts[6])
s_s.shapes.add_picture(io.BytesIO(make_png()), Inches(0.5), Inches(0.5), Inches(1), Inches(1))
prs_s.save('relswap_deck.pptx')


def _swap_rel_attrs(name, data):
    if not name.startswith('ppt/slides/_rels/'):
        return data
    xml = data.decode('utf-8')

    def swap(m):
        tag = m.group(0)
        parts = [_re.search(r'Target="[^"]*"', tag),
                 _re.search(r'TargetMode="[^"]*"', tag),
                 _re.search(r'Type="[^"]*"', tag),
                 _re.search(r'Id="[^"]*"', tag)]
        return '<Relationship ' + ' '.join(p.group(0) for p in parts if p) + '/>'

    return _re.sub(r'<Relationship [^>]*/>', swap, xml).encode('utf-8')


_rezip('relswap_deck.pptx', _swap_rel_attrs)
_b64('relswap_deck.pptx')

# SC-5 / SC-7 / SC-10 / FL-5: edgecase2_deck — astral entity, pasted-
# markdown H1 line, long digit run, malformed dcterms:modified
prs_e = Presentation()
s_e = prs_e.slides.add_slide(prs_e.slide_layouts[6])
tf = s_e.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(3)).text_frame
tf.text = "ASTRALMARK smiley survives"
tf.add_paragraph().text = "# Roadmap pasted markdown"
tf.add_paragraph().text = "digits12345678901234 glue run"
prs_e.core_properties.modified = CORE_MODIFIED
prs_e.save('edgecase2_deck.pptx')


def _edge_edits(name, data):
    if name == 'ppt/slides/slide1.xml':
        return data.replace(b'ASTRALMARK', b'&#x1F600;')
    if name == 'docProps/core.xml':
        return _re.sub(rb'(<dcterms:modified[^>]*>)[^<]*(</dcterms:modified>)',
                       rb'\1yesterday-ish\2', data)
    return data


_rezip('edgecase2_deck.pptx', _edge_edits)
_b64('edgecase2_deck.pptx')

# SC-14: encrypted_deck — edge_deck with the GP encrypted bit set on
# every local (+6) and central (+8) header
enc = bytearray(open('edge_deck.pptx', 'rb').read())
for sig, off in ((b'PK\x03\x04', 6), (b'PK\x01\x02', 8)):
    i = 0
    while True:
        j = bytes(enc).find(sig, i)
        if j < 0:
            break
        enc[j + off] |= 1
        i = j + 4
open('encrypted_deck.pptx', 'wb').write(bytes(enc))
_b64('encrypted_deck.pptx')

# SC-8: truncstored.docx — hand-built zip whose document.xml deflate
# stream is a stored block claiming 100 bytes with only 10 present
_name = b'word/document.xml'
_raw = b'\x01' + struct.pack('<HH', 100, 0xffff - 100) + b'0123456789'
_local = struct.pack('<IHHHHHIIIHH', 0x04034b50, 20, 0, 8, 0, 0, 0,
                     len(_raw), 100, len(_name), 0) + _name + _raw
_cdo = len(_local)
_central = struct.pack('<IHHHHHHIIIHHHHHII', 0x02014b50, 20, 20, 0, 8, 0, 0,
                       0, len(_raw), 100, len(_name), 0, 0, 0, 0, 0, 0) + _name
_eocd = struct.pack('<IHHHHIIH', 0x06054b50, 0, 0, 1, 1, len(_central), _cdo, 0)
open('truncstored.docx', 'wb').write(_local + _central + _eocd)
_b64('truncstored.docx')

print('batch fixtures:', {f: os.path.getsize(f) for f in
                          ('reordered_deck.pptx', 'merged_deck.pptx', 'bigimg_deck.pptx',
                           'relswap_deck.pptx', 'edgecase2_deck.pptx', 'encrypted_deck.pptx',
                           'truncstored.docx')})

# SC-14 (MediaExtract leg): encrypted deck WITH an image — MediaExtract
# only extracts media entries, so the imageless encrypted_deck never
# reaches its throw path
enc2 = bytearray(open('relswap_deck.pptx', 'rb').read())
for sig, off in ((b'PK\x03\x04', 6), (b'PK\x01\x02', 8)):
    i = 0
    while True:
        j = bytes(enc2).find(sig, i)
        if j < 0:
            break
        enc2[j + off] |= 1
        i = j + 4
open('encrypted_img_deck.pptx', 'wb').write(bytes(enc2))
_b64('encrypted_img_deck.pptx')

# ---- r2 batch fixtures (check_batch_r2.py) -------------------------------
# New-behavior fixtures for the SB-1..SB-9 batch (REVIEW_v2_5_r2.md).
# check_batch_r2.py owns their assertions until the batch is promoted.

# SB-6: hashheading_deck — content lines opening with ##/###/#### plus
# the v1.9 H1 case, in both slide body and speaker notes
prs_h = Presentation()
s_h = prs_h.slides.add_slide(prs_h.slide_layouts[5])
s_h.shapes.title.text = "HashTitle"
tf_h = s_h.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(3)).text_frame
tf_h.text = "plain line before"
tf_h.add_paragraph().text = "## Fake section pasted"
tf_h.add_paragraph().text = "### Fake notes pasted"
tf_h.add_paragraph().text = "#### deep heading pasted"
tf_h.add_paragraph().text = "# Roadmap pasted markdown"
tf_h.add_paragraph().text = "plain line after"
s_h.notes_slide.notes_text_frame.text = "## Fake heading inside notes"
prs_h.save('hashheading_deck.pptx')
_b64('hashheading_deck.pptx')

# SB-5: storednlen.docx — stored block with LEN correct and in-bounds
# but NLEN wrong (truncstored.docx covers the truncation leg)
_name2 = b'word/document.xml'
_raw2 = b'\x01' + struct.pack('<HH', 10, 0x1234) + b'0123456789'
_local2 = struct.pack('<IHHHHHIIIHH', 0x04034b50, 20, 0, 8, 0, 0, 0,
                      len(_raw2), 10, len(_name2), 0) + _name2 + _raw2
_cdo2 = len(_local2)
_central2 = struct.pack('<IHHHHHHIIIHHHHHII', 0x02014b50, 20, 20, 0, 8, 0, 0,
                        0, len(_raw2), 10, len(_name2), 0, 0, 0, 0, 0, 0) + _name2
_eocd2 = struct.pack('<IHHHHIIH', 0x06054b50, 0, 0, 1, 1, len(_central2), _cdo2, 0)
open('storednlen.docx', 'wb').write(_local2 + _central2 + _eocd2)
_b64('storednlen.docx')

# SB-8: lyingcd_deck — one referenced small png whose CENTRAL-directory
# uncompressed-size claim is patched 1000 bytes short of the truth
prs_l = Presentation()
s_l = prs_l.slides.add_slide(prs_l.slide_layouts[6])
s_l.shapes.add_picture(io.BytesIO(make_png()), Inches(0.5), Inches(0.5), Inches(1), Inches(1))
prs_l.save('lyingcd_deck.pptx')
_ly = bytearray(open('lyingcd_deck.pptx', 'rb').read())
_p = bytes(_ly).find(struct.pack('<I', 0x06054b50))
_cd = struct.unpack('<I', _ly[_p + 16:_p + 20])[0]
while _cd < _p:
    assert struct.unpack('<I', _ly[_cd:_cd + 4])[0] == 0x02014b50
    _nl, _xl, _cl = struct.unpack('<HHH', _ly[_cd + 28:_cd + 34])
    _nm = bytes(_ly[_cd + 46:_cd + 46 + _nl])
    if _nm.startswith(b'ppt/media/') and _nm.endswith(b'.png'):
        _usz = struct.unpack('<I', _ly[_cd + 24:_cd + 28])[0]
        _ly[_cd + 24:_cd + 28] = struct.pack('<I', _usz - 1000)
    _cd += 46 + _nl + _xl + _cl
open('lyingcd_deck.pptx', 'wb').write(bytes(_ly))
_b64('lyingcd_deck.pptx')

# SB-7: manytables.docx — hand-assembled document.xml with 205 tiny
# tables (past the 200-table rendering guard)
_tbls = ''.join('<w:tbl><w:tr><w:tc><w:p><w:r><w:t>tbl%dcell</w:t></w:r></w:p></w:tc></w:tr></w:tbl>' % i
                for i in range(205))
_doc = ('<?xml version="1.0"?><w:document><w:body>' + _tbls +
        '<w:p><w:r><w:t>after the tables</w:t></w:r></w:p></w:body></w:document>')
with zipfile.ZipFile('manytables.docx', 'w', zipfile.ZIP_DEFLATED) as _z:
    _z.writestr('word/document.xml', _doc)
    _z.writestr('word/_rels/document.xml.rels',
                '<?xml version="1.0"?><Relationships/>')
_b64('manytables.docx')

print('r2 fixtures:', {f: os.path.getsize(f) for f in
                       ('hashheading_deck.pptx', 'storednlen.docx',
                        'lyingcd_deck.pptx', 'manytables.docx')})

# ---- r6 batch fixtures (check_batch_r6.py) -------------------------------
# CF-1: code_deck — an Arcade script pasted across slide paragraphs
# (internal blank line, a '# ...' comment line that SB-6 will escape,
# indented lines), code-shaped level-1 bullets, and prose/table control
# content that must never be absorbed into a fence.

prs_c = Presentation()
s_c = prs_c.slides.add_slide(prs_c.slide_layouts[5])
s_c.shapes.title.text = "CodeTitle Arcade Sample"
tf_c = s_c.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(4)).text_frame
tf_c.text = "Input expression for the stationing calculation below:"
for txt in ["var station = $feature.MEASURE",
            "var IsNegative = 0 //tracks values",
            " ",  # space-only run -> a blank line inside the script
            "if (station < 0) {",
            "  station = station * -1",
            "}",
            "# stationing format note",
            "return Stationlb",
            "Test with UNAPR data and Pipeline Referencing centerlines today"]:
    tf_c.add_paragraph().text = txt
p_c = tf_c.add_paragraph()
p_c.text = "$feature.Depth + $feature.Width / 10"
p_c.level = 1
p_c = tf_c.add_paragraph()
p_c.text = "positive case expected result route measure"
p_c.level = 1
tbl_c = s_c.shapes.add_table(2, 2, Inches(0.5), Inches(6), Inches(6), Inches(1)).table
tbl_c.cell(0, 0).text = "a = b; c = d;"
tbl_c.cell(0, 1).text = "plain cell"
tbl_c.cell(1, 0).text = "var x = 1;"
tbl_c.cell(1, 1).text = "referent"
prs_c.save('code_deck.pptx')
_b64('code_deck.pptx')

# CF-1 control deck: prose that superficially flirts with code shapes —
# instruction lines ending in ';', a lowercase 'return to ...' sentence
# — must come through byte-identically to a fence-free extraction
prs_p = Presentation()
s_p = prs_p.slides.add_slide(prs_p.slide_layouts[5])
s_p.shapes.title.text = "ProseTitle"
tf_p = s_p.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(3)).text_frame
tf_p.text = "Select the route; click Save; verify the label renders;"
tf_p.add_paragraph().text = "return to the map view and verify the resulting label"
tf_p.add_paragraph().text = "Numerous scripts can be found at https://github.com/Esri/arcade-expressions today"
prs_p.save('prose_deck.pptx')
_b64('prose_deck.pptx')

print('r6 fixtures:', {f: os.path.getsize(f) for f in
                       ('code_deck.pptx', 'prose_deck.pptx')})

# ---- v2.2 fixtures (DL-1 diagram-label collapse) --------------------------
# diagram_deck: slide 1 draws a route diagram the corpus way — dashed
# connectors plus tiny floating text boxes for tick numbers and
# route/event ids — alongside real content (a prose paragraph, a table,
# a long floating callout that must stay inline). Slide 2 is the
# control: only 3 short floating boxes, below the cluster threshold, so
# they must stay inline and no [figure: ...] line may appear.
# Owned by check_format.py §12; not in the recall set (tick numbers
# deliberately compress to "10–15").
from pptx.enum.shapes import MSO_CONNECTOR

prs_d = Presentation()
s_d = prs_d.slides.add_slide(prs_d.slide_layouts[5])
s_d.shapes.title.text = "DiagramTitle Route Split"
tf_d = s_d.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(1)).text_frame
tf_d.text = "Verify the split renders across the diagram below today"
for k, lbl in enumerate(['10', '11', '12', '13', '14', '15',
                         'R1', 'E1', 'E1', 'Output']):
    s_d.shapes.add_textbox(Inches(0.4 + 0.55 * k), Inches(2.6),
                           Inches(0.5), Inches(0.3)).text_frame.text = lbl
s_d.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.4), Inches(3.2),
                         Inches(6.0), Inches(3.2))
s_d.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(6), Inches(0.5)) \
    .text_frame.text = "Modify this test case to use measure five instead"
tbl_d = s_d.shapes.add_table(2, 2, Inches(0.5), Inches(4.5), Inches(6), Inches(1)).table
tbl_d.cell(0, 0).text = 'Route ID'
tbl_d.cell(0, 1).text = 'R1'
tbl_d.cell(1, 0).text = 'Measure'
tbl_d.cell(1, 1).text = '10'
s_d2 = prs_d.slides.add_slide(prs_d.slide_layouts[5])
s_d2.shapes.title.text = "ControlTitle"
for k, lbl in enumerate(['A1', 'B2', 'C3']):
    s_d2.shapes.add_textbox(Inches(0.5 + k), Inches(2), Inches(0.6), Inches(0.3)) \
        .text_frame.text = lbl
prs_d.save('diagram_deck.pptx')
_b64('diagram_deck.pptx')

print('v2.2 fixtures:', {f: os.path.getsize(f) for f in ('diagram_deck.pptx',)})


# SB-5 (MediaExtract leg): storednlen_img.pptx — a single media entry
# whose deflate stream is a stored block with a wrong NLEN. MediaExtract
# only inflates media entries, so the docx variant never reaches its
# inflate; this one does. (Not a real png — the zip reader never checks.)
_name3 = b'ppt/media/image1.png'
_raw3 = b'\x01' + struct.pack('<HH', 10, 0x1234) + b'0123456789'
_local3 = struct.pack('<IHHHHHIIIHH', 0x04034b50, 20, 0, 8, 0, 0, 0,
                      len(_raw3), 10, len(_name3), 0) + _name3 + _raw3
_cdo3 = len(_local3)
_central3 = struct.pack('<IHHHHHHIIIHHHHHII', 0x02014b50, 20, 20, 0, 8, 0, 0,
                        0, len(_raw3), 10, len(_name3), 0, 0, 0, 0, 0, 0) + _name3
_eocd3 = struct.pack('<IHHHHIIH', 0x06054b50, 0, 0, 1, 1, len(_central3), _cdo3, 0)
open('storednlen_img.pptx', 'wb').write(_local3 + _central3 + _eocd3)
_b64('storednlen_img.pptx')


# ---- v2.5 (Sidecar_Format_Plan phase 2): structure fixtures ---------------
# cells_deck.pptx — CP-1 cell paragraphs (single-column label box →
# bold label + bullets; multi-column cell → <br>), TP-1 topmost label
# of a title-less slide becomes the heading even though it sits AFTER
# the table in z-order, IB-1 inherited body-placeholder bullets (with
# one explicit buNone opt-out).
from pptx.oxml.ns import qn as _qn
from lxml import etree as _et
prs_v = Presentation()
s_v1 = prs_v.slides.add_slide(prs_v.slide_layouts[6])            # blank: no title placeholder
_t1 = s_v1.shapes.add_table(2, 1, Inches(0.5), Inches(1.5), Inches(6), Inches(2)).table
_t1.cell(0, 0).text = 'Positive Tests: Normal Routes'
_tf = _t1.cell(1, 0).text_frame
_tf.text = 'Correct line order of 100, 200, 300, 400 on a normal line'
for _txt in ('Correct line order of 300, 400, 500, 600 on a normal line',
             'Time sliced routes, first slice 100, 200 and second 300, 400'):
    _tf.add_paragraph().text = _txt
_t2 = s_v1.shapes.add_table(2, 3, Inches(0.5), Inches(4), Inches(6), Inches(1.5)).table
for _c, _h in enumerate(('#', 'Test', 'Expected result')):
    _t2.cell(0, _c).text = _h
_t2.cell(1, 0).text = 'A-1'
_t2.cell(1, 1).text = 'Toggle is present'
_tf2 = _t2.cell(1, 2).text_frame
_tf2.text = 'Toggle shown'
_tf2.add_paragraph().text = 'Default OFF'
# the section label: drawn last (top of z-order) but positioned at the top
s_v1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.6)).text_frame.text = 'Coordinate Configuration Tests'
s_v1.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(6), Inches(0.6)).text_frame.text = 'Footnote below the tables'
# v2.6 (TP-2): a title-less slide whose TOPMOST short text is a route
# label of the drawn diagram (four labels — a cluster) must NOT take
# that label as its heading; its case line stays in the body
s_v3 = prs_v.slides.add_slide(prs_v.slide_layouts[6])
for _k, _lbl in enumerate(('1A_New; 100', '1B; 200', '1C; 300', '2A, 100')):
    s_v3.shapes.add_textbox(Inches(0.5 + 1.5 * _k), Inches(0.3), Inches(1.3), Inches(0.4)).text_frame.text = _lbl
s_v3.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8), Inches(0.6)).text_frame.text = (
    'Reassign all the routes in a line to another line on right, transferring routes.')
s_v3.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(8), Inches(0.6)).text_frame.text = (
    '2: Transfer to an existing line – spanning Events – Stayput and Retire Behavior.')
s_v2 = prs_v.slides.add_slide(prs_v.slide_layouts[1])            # title + content
s_v2.shapes.title.text = 'Acceptance Criteria'
_body = s_v2.placeholders[1].text_frame
_body.text = 'Shows only for an UN-APR dataset'
_body.add_paragraph().text = 'Allow turning these layers ON and OFF'
_opt = _body.add_paragraph()
_opt.text = 'Plain note with the bullet turned off'
_ppr = _opt._p.get_or_add_pPr()
_ppr.append(_et.SubElement(_ppr, _qn('a:buNone')))
prs_v.save('cells_deck.pptx')
_b64('cells_deck.pptx')

# labels.docx — DL-2: bold label paragraphs and ":"-labels followed by a
# list become "### " headings; List Number renders as "1. "; a long
# bold sentence and a plain paragraph stay as they are.
doc_l = Document()
doc_l.add_paragraph().add_run('UI Tests – First Pane').bold = True
doc_l.add_paragraph('Verify the Open Type is set from the configuration', style='List Bullet')
doc_l.add_paragraph('Verify the Attribute Set is as configured', style='List Bullet')
doc_l.add_paragraph('Negative Tests:')
doc_l.add_paragraph('Verify Next is disabled when a required field is empty', style='List Bullet')
doc_l.add_paragraph('Steps to reproduce:')
doc_l.add_paragraph('Open the widget', style='List Number')
doc_l.add_paragraph('Pick a route', style='List Number')
doc_l.add_paragraph().add_run('This whole sentence is bold but far too long to be a section label in any document.').bold = True
doc_l.add_paragraph('A plain paragraph that ends with a colon but has no list after it:')
doc_l.add_paragraph('Just prose here.')
doc_l.save('labels.docx')
_b64('labels.docx')
print('v2.5 fixtures:', {f: os.path.getsize(f) for f in ('cells_deck.pptx', 'labels.docx')})
