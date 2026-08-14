"""Standing media suite for the current scripts/ generation.

Since ZipTextExtract v2.2 (MG-1, the r7 batch) media extraction lives
inside ZipTextExtract itself: the optional fourth parameter `wantMedia`
makes the SC-4 selection walk also inflate and return the selected
rasters (`images` / `skippedMedia` / `imageCount`). MediaExtract is
RETIRED with the flow v2.9 window, so this suite — not the MediaExtract
legs inside check_format.py — is where media coverage lives from r7 on.

Contracts:

  1. selection: caps (12 / 350 KB / ~3 MB, central-directory order)
     decide both the minted links and the returned images — the
     under-cap image is linked AND returned, the over-cap one is
     neither (it lands in skippedMedia)
  2. SC-4 structural: every image link minted in `text` names an entry
     in `images`, and `media` lists exactly the linked names
  3. byte fidelity: each returned b64 decodes to the exact bytes of the
     source archive's media entry
  4. off switch: wantMedia absent -> images == [] / imageCount == 0,
     and every other field byte-identical to the wantMedia run (the
     walk may not change what it selects when it also saves)
  5. throw paths (wantMedia only): lying central-directory size claim,
     stored-block NLEN corruption on a media entry, encrypted entries.
     WITHOUT wantMedia the lying archive must NOT throw — media entries
     are never inflated, exactly as pre-v2.2
  6. budget: a worst-case archive (12 near-cap images + text) keeps the
     serialized result under 4.7 MB — inside the ~5 MB Run-script
     response cap with margin

Prereqs: make_fixtures.py has run in this directory.
Exit code: non-zero on any failed assertion.
"""
import base64
import io
import json
import os
import re
import struct
import subprocess
import sys
import zipfile
import zlib

SCRIPTS = os.environ.get('HARNESS_SCRIPTS', '../../scripts')

failures = []


def check(cond, label):
    print(('ok   ' if cond else 'FAIL ') + label)
    if not cond:
        failures.append(label)


# ---- wrapped runner: same shape as wrap.py plus the wantMedia flag ------
src = open(f'{SCRIPTS}/ZipTextExtract.ts', encoding='utf-8').read().replace(
    'workbook: ExcelScript.Workbook', 'workbook: unknown')
src += '''
// ---- harness appendix ----
declare const require: (m: string) => { readFileSync: (p: string, e: string) => string };
const fs = require('fs');
const argv = (globalThis as {process?: {argv: string[]}}).process!.argv;
const b64 = fs.readFileSync(argv[2], 'utf8');
const label = argv[3] || '';
const want = argv[4] === 'media';
console.log(JSON.stringify({out: main(null as unknown, b64, label, want)}));
'''
open('zte_media.ts', 'w', encoding='utf-8').write(src)


def run(fixture, prefix='media/docX_', want=True):
    return subprocess.run(['node', '--experimental-strip-types', 'zte_media.ts',
                           fixture, prefix, 'media' if want else ''],
                          capture_output=True, text=True, encoding='utf-8')


def out(r):
    return json.loads(r.stdout)['out']


# ---- 1: cap-aware selection on bigimg_deck ------------------------------
with zipfile.ZipFile('bigimg_deck.pptx') as z:
    sizes = {n.split('/')[-1]: z.getinfo(n).file_size for n in z.namelist() if '/media/' in n}
    member = {n.split('/')[-1]: z.read(n) for n in z.namelist() if '/media/' in n}
small = [n for n, sz in sizes.items() if sz <= 350 * 1024][0]
big = [n for n, sz in sizes.items() if sz > 350 * 1024][0]
b = out(run('bigimg_deck.pptx.b64'))
check(b['imageCount'] == 1 and [i['name'] for i in b['images']] == [small],
      'bigimg: only the under-cap image is returned')
check(big in b['skippedMedia'].split('\n') and small not in b['skippedMedia'],
      'bigimg: the over-cap image lands in skippedMedia')
check(f'](media/docX_{small})' in b['text'] and big not in b['text'],
      'bigimg: link set matches the returned set')

# ---- 3: byte fidelity ---------------------------------------------------
check(base64.b64decode(b['images'][0]['b64']) == member[small],
      'bigimg: returned b64 decodes to the exact source bytes')

# ---- 2: SC-4 structural on the real deck --------------------------------
r = out(run('real_deck.pptx.b64'))
linked = set(re.findall(r'!\[[^\]]*\]\(media/docX_([^)]+)\)', r['text']))
returned = {i['name'] for i in r['images']}
check(r['imageCount'] >= 1 and linked and linked <= returned,
      f'real_deck: every minted link names a returned image '
      f'({len(linked)} linked, {r["imageCount"]} returned)')
check(set(r['media'].split('\n')) == linked,
      'real_deck: media lists exactly the linked names')
with zipfile.ZipFile('real_deck.pptx') as z:
    rmember = {n.split('/')[-1]: z.read(n) for n in z.namelist() if '/media/' in n}
check(all(base64.b64decode(i['b64']) == rmember[i['name']] for i in r['images']),
      'real_deck: every returned image is byte-exact')

# ---- 4: the off switch is a strict no-op --------------------------------
plain = out(run('real_deck.pptx.b64', want=False))
check(plain['images'] == [] and plain['imageCount'] == 0,
      'wantMedia absent: images empty, imageCount 0')
same = all(plain[k] == r[k] for k in
           ('text', 'rels', 'parts', 'kind', 'media', 'author',
            'lastEditedBy', 'lastEdited', 'skippedMedia'))
check(same, 'wantMedia absent: every other field identical to the media run')

# ---- 5: throw paths -----------------------------------------------------
f = run('lyingcd_deck.pptx.b64')
check(f.returncode != 0 and 'size mismatch' in f.stderr,
      'lying central directory throws under wantMedia')
f = run('lyingcd_deck.pptx.b64', want=False)
check(f.returncode == 0,
      'lying central directory does NOT throw without wantMedia '
      '(media never inflated, the pre-v2.2 behavior)')
f = run('storednlen_img.pptx.b64', prefix='')
check(f.returncode != 0 and 'NLEN mismatch' in f.stderr,
      'media stored-block NLEN corruption throws under wantMedia')
f = run('encrypted_img_deck.pptx.b64', prefix='')
check(f.returncode != 0 and 'encrypted' in f.stderr,
      'encrypted media entry throws under wantMedia')

# ---- 6: budget — worst-case payload stays under the response cap --------
# 12 near-cap random-data PNGs (~250 KB each, incompressible) + text:
# media b64 lands just under 4.1 MB, the true worst case the 12 / 350 KB
# / 3 MB caps allow.
BUDGET = 'budget_deck.pptx'
if not os.path.exists(BUDGET + '.b64'):
    from pptx import Presentation
    from pptx.util import Inches

    def make_png(w, h, seed):
        # distinct bytes per image — python-pptx dedups identical
        # pictures into one media part, which would defeat the leg
        import random
        random.seed(seed)

        def chunk(tag, data):
            c = tag + data
            return (struct.pack('>I', len(data)) + c
                    + struct.pack('>I', zlib.crc32(c) & 0xffffffff))
        raw = b''.join(b'\x00' + bytes(random.randrange(256) for _ in range(w * 3))
                       for _ in range(h))
        return (b'\x89PNG\r\n\x1a\n'
                + chunk(b'IHDR', struct.pack('>IIBBBBB', w, h, 8, 2, 0, 0, 0))
                + chunk(b'IDAT', zlib.compress(raw))
                + chunk(b'IEND', b''))

    prs = Presentation()
    for i in range(12):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f'Budget slide {i}'
        tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(3)).text_frame
        tb.text = ('route measure referent calibration ' * 60).strip()
        # ~256 KB each; 12 of these fit under MAX_TOTAL
        s.shapes.add_picture(io.BytesIO(make_png(292, 292, i)), Inches(8.6),
                             Inches(0.2), Inches(0.6), Inches(0.6))
    prs.save(BUDGET)
    with open(BUDGET, 'rb') as fh:
        open(BUDGET + '.b64', 'w', encoding='utf-8').write(
            base64.b64encode(fh.read()).decode())

bres = run(BUDGET + '.b64')
bout = out(bres)
payload = len(json.dumps(bout))
check(bout['imageCount'] >= 10,
      f'budget: near-cap image set actually saved ({bout["imageCount"]} images)')
check(payload > 3_000_000,
      f'budget: the leg exercises a genuinely large payload ({payload} bytes)')
check(payload < 4_700_000,
      f'budget: worst-case result stays under 4.7 MB ({payload} bytes)')

print()
if failures:
    print(f'RESULT: FAIL — {len(failures)} assertion(s) failed')
    sys.exit(1)
print('RESULT: PASS')
