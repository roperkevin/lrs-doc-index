"""sidecar_to_pptx — render a user-story sidecar .md as a designed .pptx briefing deck.

    python3 pptxgen/sidecar_to_pptx.py INPUT.md [-o OUT.pptx] [--media-dir DIR]
            [--enhance] [--model MODEL] [--strict-enhance]
            [--outline-json PATH] [--from-outline]

Pipeline:  sidecar_parser -> DeckOutline (deck_outline) -> this renderer.
The deterministic outline mapping is the default; --enhance swaps in the
Claude API stage (pptxgen/enhance.py) which restructures the extracted
text into a presentation-grade outline. Enhancement failures fall back
to the deterministic outline with a warning unless --strict-enhance.

Exit codes: 0 success; 2 input is not a sidecar (no ```yaml fence);
3 --strict-enhance and the enhancement stage failed.

Design: "Midnight Executive" — navy-dominant sandwich (dark title and
closers, light content), Cambria headers / Calibri body (safe-list
fonts), rounded-rect chip/card motif throughout. python-pptx cannot
measure text, so layout is estimate-based (average Calibri advance
width) with overflow handled by font stepping and greedy pagination —
long sections continue on "(cont.)" slides, tables repeat their header
row per part, and nothing renders past the content region.

Gate: review/harness/check_pptx.py (wired into .github/workflows/
harness.yml full-format). Visual QA is a documented local step — see
pptxgen/README.md.
"""
from __future__ import annotations

import argparse
import json
import math
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck_outline import (DeckOutline, OBlock, OSlide,  # noqa: E402
                          outline_from_sidecar, load_outline)
from sidecar_parser import SidecarError, parse_sidecar  # noqa: E402

# ---------------------------------------------------------------- theme

NAVY = RGBColor(0x1E, 0x27, 0x61)       # dominant
NAVY_2 = RGBColor(0x2A, 0x37, 0x7E)     # cards on dark slides / callouts
ICE = RGBColor(0xCA, 0xDC, 0xFC)        # chips, accents on dark
ICE_BG = RGBColor(0xEE, 0xF3, 0xFC)     # tints, banding, cards on light
LIGHT_ICE = RGBColor(0xD8, 0xE2, 0xF5)  # secondary text on dark cards
ICE_MUTED = RGBColor(0x8F, 0xA6, 0xD9)  # muted text on dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x20, 0x26, 0x33)        # body text on light
MUTED = RGBColor(0x5A, 0x64, 0x78)      # captions on light

FONT_HEAD = 'Cambria'
FONT_BODY = 'Calibri'

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
MARGIN = 0.6
BODY_X, BODY_Y = MARGIN, 1.52
BODY_W = SLIDE_W_IN - 2 * MARGIN        # 12.133
BODY_H = 5.45                           # bottom edge at 6.97
TABLE_BODY_ROWS_MAX = 8                 # body rows per slide part

# ------------------------------------------------------------ estimates
# python-pptx cannot measure text; these are the calibrated estimates
# the pagination is built on. 0.52em average advance for Calibri, 1.28
# line spacing factor.


def _chars_per_line(font_pt: float, width_in: float) -> int:
    # 72 points per inch; 0.52 em average advance width for Calibri
    return max(8, int(width_in * 72 / (font_pt * 0.52)))


def _text_lines(text: str, font_pt: float, width_in: float) -> int:
    cpl = _chars_per_line(font_pt, width_in)
    return max(1, sum(max(1, math.ceil(len(part) / cpl))
                      for part in (text.split('\n') or [''])))


def _text_height(text: str, font_pt: float, width_in: float,
                 space_after: float = 0.07) -> float:
    return _text_lines(text, font_pt, width_in) * font_pt * 1.28 / 72 \
        + space_after


def _table_col_widths(rows: list[list[str]], max_w: float) -> list[float]:
    n_cols = max(len(r) for r in rows)
    weights = []
    for c in range(n_cols):
        longest = max((len(r[c]) for r in rows if c < len(r)), default=1)
        weights.append(min(4.5, max(0.9, longest * 0.09)))
    scale = min(1.0, max_w / sum(weights))
    return [w * scale for w in weights]


def _table_row_heights(rows: list[list[str]], widths: list[float],
                       font_pt: float) -> list[float]:
    heights = []
    for r in rows:
        lines = 1
        for c, cell in enumerate(r):
            if c < len(widths):
                lines = max(lines, _text_lines(cell, font_pt,
                                               widths[c] - 0.16))
        heights.append(max(0.32, lines * font_pt * 1.2 / 72 + 0.1))
    return heights


# -------------------------------------------------------------- helpers

def _set_font(run, size, color, bold=False, italic=False, font=FONT_BODY):
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic


def _textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, margin=0.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, Inches(margin))
    return tb, tf


def _para(tf, first: bool):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def _set_indent(p, mar_l_in: float, hang_in: float = 0.0) -> None:
    ppr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    ppr.set('marL', str(int(Inches(mar_l_in))))
    if hang_in:
        ppr.set('indent', str(-int(Inches(hang_in))))


def _rounded_rect(slide, x, y, w, h, fill, *, radius=0.25, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _chip(slide, x, y, text, *, fill=ICE_BG, color=MUTED, size=11,
          bold=False, pad=0.14, h=0.32, caps=False):
    label = text.upper() if caps else text
    w = pad * 2 + len(label) * size * 0.62 / 72
    shp = _rounded_rect(slide, x, y, w, h, fill, radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, Inches(0.02))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    _set_font(r, size, color, bold=bold)
    return w


def _dark_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY


def _light_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def _meta_str(meta: dict, key: str) -> str:
    v = meta.get(key)
    return '' if v is None else str(v).strip()


def _fmt_date(iso: str) -> str:
    return iso[:10] if iso else ''


# ------------------------------------------------------------- renderer

class DeckRenderer:
    def __init__(self, outline: DeckOutline, media_dir: str | None,
                 base_dir: str):
        self.o = outline
        self.media_dir = media_dir
        self.base_dir = base_dir
        self.prs = Presentation()
        self.prs.slide_width = Emu(12192000)   # 13.333 in
        self.prs.slide_height = Emu(6858000)   # 7.5 in
        self.blank = self.prs.slide_layouts[6]

    # -- top-level ----------------------------------------------------
    def build(self) -> Presentation:
        for oslide in self.o.slides:
            kind = oslide.kind
            if kind == 'title':
                self._title_slide()
            elif kind == 'summary':
                self._summary_slide()
            elif kind == 'agenda':
                self._agenda_slide(oslide)
            elif kind == 'content':
                self._content_slides(oslide)
            elif kind == 'takeaways':
                self._takeaways_slide(oslide)
            elif kind == 'related':
                self._related_slides()
        return self.prs

    def _new_slide(self, dark: bool):
        slide = self.prs.slides.add_slide(self.blank)
        (_dark_background if dark else _light_background)(slide)
        return slide

    # -- title slide ---------------------------------------------------
    def _title_slide(self) -> None:
        slide = self._new_slide(dark=True)
        meta = self.o.meta
        kind = _meta_str(meta, 'doc_kind') or 'Document'
        _chip(slide, MARGIN, 0.85, kind, fill=ICE, color=NAVY, size=12,
              bold=True, h=0.38, caps=True)

        title = self.o.title or _meta_str(meta, 'title') or 'Untitled'
        size = 40
        while size > 28 and _text_lines(title, size, BODY_W) > 3:
            size -= 6
        _, tf = _textbox(slide, MARGIN, 1.6, BODY_W, 2.6)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_font(r, size, WHITE, bold=True, font=FONT_HEAD)

        pairs = [('Surface', _meta_str(meta, 'surface')),
                 ('Target release', _meta_str(meta, 'target_release')),
                 ('PE', _meta_str(meta, 'pe')),
                 ('Dev', _meta_str(meta, 'dev'))]
        parts = [f'{k}: {v}' for k, v in pairs if v]
        if parts:
            _, tf = _textbox(slide, MARGIN, 4.5, BODY_W, 0.5)
            p = tf.paragraphs[0]
            for j, part in enumerate(parts):
                if j:
                    r = p.add_run()
                    r.text = '   ·   '
                    _set_font(r, 14, ICE_MUTED)
                label, _, value = part.partition(': ')
                r = p.add_run()
                r.text = f'{label}  '
                _set_font(r, 14, ICE_MUTED)
                r = p.add_run()
                r.text = value
                _set_font(r, 14, ICE, bold=True)

        strip = ' · '.join(x for x in (
            self.o.source_file,
            f"Extracted {_fmt_date(_meta_str(meta, 'extracted'))}"
            if _meta_str(meta, 'extracted') else '',
            _meta_str(meta, 'extraction_lane')) if x)
        if strip:
            _, tf = _textbox(slide, MARGIN, 6.55, BODY_W, 0.4)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = strip
            _set_font(r, 11, ICE_MUTED)

    # -- summary slide ---------------------------------------------------
    def _summary_slide(self) -> None:
        slide = self._new_slide(dark=False)
        self._slide_title(slide, 'Summary')

        left_w = 7.7
        summary = self.o.summary or '(No summary available.)'
        size = 20
        while size > 15 and _text_height(summary, size, left_w - 0.5) > 3.1:
            size -= 2

        # oversized open-quote glyph as the visual anchor
        _, tf = _textbox(slide, MARGIN - 0.08, 1.28, 1.0, 1.2)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = '“'
        _set_font(r, 88, ICE, bold=True, font=FONT_HEAD)

        _, tf = _textbox(slide, MARGIN + 0.55, 1.62, left_w - 0.55, 3.3)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = summary
        _set_font(r, size, INK)
        p.line_spacing = 1.15

        # keyword chips, max two rows
        if self.o.keywords:
            x, y = MARGIN, 5.35
            rows = 0
            shown = 0
            for kw in self.o.keywords:
                w = 0.28 + len(kw) * 12 * 0.62 / 72
                if x + w > MARGIN + left_w:
                    rows += 1
                    if rows >= 2:
                        break
                    x, y = MARGIN, y + 0.46
                _chip(slide, x, y, kw, fill=ICE_BG, color=NAVY, size=12,
                      h=0.36)
                x += w + 0.14
                shown += 1
            if shown < len(self.o.keywords):
                _chip(slide, x, y, f'+{len(self.o.keywords) - shown} more',
                      fill=ICE_BG, color=MUTED, size=12, h=0.36)

        self._metadata_card(slide, x=8.75, y=1.42, w=3.98, h=5.2)

    def _metadata_card(self, slide, x, y, w, h) -> None:
        _rounded_rect(slide, x, y, w, h, ICE_BG, radius=0.06)
        meta = self.o.meta
        fields = [('Doc ID', _meta_str(meta, 'doc_id')),
                  ('Target release', _meta_str(meta, 'target_release')),
                  ('PE', _meta_str(meta, 'pe')),
                  ('Dev', _meta_str(meta, 'dev')),
                  ('Last edited',
                   ' '.join(p for p in
                            (_fmt_date(_meta_str(meta, 'last_edited')),
                             ('by ' + _meta_str(meta, 'last_edited_by'))
                             if _meta_str(meta, 'last_edited_by') else '')
                            if p)),
                  ('Extracted', _fmt_date(_meta_str(meta, 'extracted'))),
                  ('Lane', _meta_str(meta, 'extraction_lane')),
                  ('Prompt version', _meta_str(meta, 'prompt_version'))]
        fields = [(k, v) for k, v in fields if v]
        _, tf = _textbox(slide, x + 0.3, y + 0.25, w - 0.6, h - 0.5)
        first = True
        for label, value in fields:
            p = _para(tf, first)
            first = False
            r = p.add_run()
            r.text = label.upper()
            _set_font(r, 10, MUTED, bold=True)
            p.space_before = Pt(0 if p is tf.paragraphs[0] else 9)
            p2 = tf.add_paragraph()
            r = p2.add_run()
            r.text = value
            _set_font(r, 13, INK)

    # -- agenda slide ---------------------------------------------------
    def _agenda_slide(self, oslide: OSlide) -> None:
        slide = self._new_slide(dark=False)
        self._slide_title(slide, oslide.title or 'Agenda')
        items = oslide.items
        per_col = max(1, math.ceil(len(items) / (2 if len(items) > 6 else 1)))
        col_w = (BODY_W - 0.6) / 2 if len(items) > 6 else BODY_W
        row_h = min(0.78, BODY_H / per_col)
        for idx, item in enumerate(items):
            col, row = divmod(idx, per_col)
            x = BODY_X + col * (col_w + 0.6)
            y = BODY_Y + 0.1 + row * row_h
            badge = _rounded_rect(slide, x, y, 0.42, 0.42, ICE, radius=0.5)
            btf = badge.text_frame
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            btf.word_wrap = False
            for m in ('margin_left', 'margin_right', 'margin_top',
                      'margin_bottom'):
                setattr(btf, m, Inches(0.0))
            p = btf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(idx + 1)
            _set_font(r, 14, NAVY, bold=True)
            _, tf = _textbox(slide, x + 0.62, y - 0.02, col_w - 0.62, row_h,
                             anchor=MSO_ANCHOR.TOP)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = item
            _set_font(r, 15, INK)

    # -- content slides (with pagination) --------------------------------
    @staticmethod
    def _title_metrics(text: str, kicker: str, subhead: str) -> tuple[int, float]:
        """(font size, bottom edge in inches) for a slide title block."""
        size = 28
        title_w = BODY_W - (1.6 if kicker else 0)
        while size > 20 and _text_lines(text, size, title_w) > 2:
            size -= 4
        lines = _text_lines(text, size, title_w)
        bottom = 0.42 + lines * size * 1.18 / 72
        if subhead:
            bottom += _text_lines(subhead, 15, title_w) * 15 * 1.25 / 72 + 0.06
        return size, bottom

    def _slide_title(self, slide, text, *, subhead='', kicker='',
                     dark=False) -> None:
        color = WHITE if dark else NAVY
        size, _ = self._title_metrics(text, kicker, subhead)
        title_w = BODY_W - (1.6 if kicker else 0)
        _, tf = _textbox(slide, MARGIN, 0.42, title_w, 1.6)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        _set_font(r, size, color, bold=True, font=FONT_HEAD)
        if subhead:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = subhead
            _set_font(r, 15, ICE_MUTED if dark else MUTED)
            p.space_before = Pt(4)
        if kicker:
            _chip(slide, SLIDE_W_IN - MARGIN - 1.35, 0.5, kicker,
                  fill=ICE_BG, color=MUTED, size=11, h=0.32)

    def _block_height(self, b: OBlock, font_pt: float) -> float:
        if b.kind in ('para', 'persona'):
            width = BODY_W - (1.1 if b.kind == 'persona' else 0)
            h = _text_height(b.text, font_pt, width)
            return h + 0.45 if b.kind == 'persona' else h
        if b.kind == 'bullet':
            indent = 0.28 * min(b.level, 4)
            return _text_height(b.text, font_pt, BODY_W - indent - 0.3)
        if b.kind == 'subheading':
            return _text_height(b.text, 18 if b.level == 3 else 15, BODY_W,
                                space_after=0.12)
        if b.kind == 'table':
            rows = b.rows[:TABLE_BODY_ROWS_MAX + 1]
            widths = _table_col_widths(rows, BODY_W)
            font_size = 10 if any(len(c) > 120 for r in rows for c in r) else 12
            return sum(_table_row_heights(rows, widths, font_size)) + 0.15
        if b.kind == 'image':
            return self._image_height(b) + 0.12
        return 0.3

    def _image_height(self, b: OBlock) -> float:
        path = self._resolve_media(b.path)
        max_w, max_h = 7.5, 3.2
        if path:
            try:
                from PIL import Image as PILImage
                with PILImage.open(path) as im:
                    w, h = im.size
                scale = min(max_w / (w / 96), max_h / (h / 96), 1.0)
                return max(0.6, (h / 96) * scale)
            except Exception:
                pass
        return 1.4  # placeholder height

    def _resolve_media(self, link: str) -> str | None:
        if not link:
            return None
        name = os.path.basename(link.replace('\\', '/'))
        if self.media_dir:
            cand = os.path.join(self.media_dir, name)
        else:
            cand = os.path.normpath(os.path.join(self.base_dir, link))
        return cand if os.path.isfile(cand) else None

    def _paginate(self, oslide: OSlide,
                  avail_h: float) -> list[tuple[list[OBlock], float]]:
        """Split a content slide's blocks into pages of the body region.
        Returns [(blocks, font_pt)] per page."""
        blocks = oslide.blocks
        if not blocks:
            return [([], 16.0)]

        font = 16.0
        total = sum(self._block_height(b, font) for b in blocks)
        simple = all(b.kind in ('para', 'bullet') for b in blocks)
        if total > avail_h and simple and total <= avail_h * 1.35:
            for step in (14.0, 12.0):
                font = step
                total = sum(self._block_height(b, font) for b in blocks)
                if total <= avail_h:
                    break

        pages: list[tuple[list[OBlock], float]] = []
        cur: list[OBlock] = []
        used = 0.0

        def close() -> None:
            nonlocal cur, used
            if cur:
                pages.append((cur, font))
            cur, used = [], 0.0

        for b in blocks:
            # split oversize tables into <=8-body-row parts, header repeated
            if b.kind == 'table' and len(b.rows) - 1 > TABLE_BODY_ROWS_MAX:
                header, body = b.rows[0], b.rows[1:]
                for i in range(0, len(body), TABLE_BODY_ROWS_MAX):
                    part = OBlock(kind='table',
                                  rows=[header] + body[i:i + TABLE_BODY_ROWS_MAX])
                    h = self._block_height(part, font)
                    if cur and used + h > avail_h:
                        close()
                    cur.append(part)
                    used += h
                continue
            h = self._block_height(b, font)
            if cur and used + h > avail_h:
                close()
            cur.append(b)
            used += min(h, avail_h)  # a single oversize block gets its own page
        close()
        return pages or [([], font)]

    def _content_slides(self, oslide: OSlide) -> None:
        _, title_bottom = self._title_metrics(oslide.title + ' (cont.)',
                                              oslide.kicker, oslide.subhead)
        body_y = max(BODY_Y, title_bottom + 0.12)
        avail_h = 6.97 - body_y
        pages = self._paginate(oslide, avail_h)
        for page_no, (blocks, font) in enumerate(pages):
            slide = self._new_slide(dark=False)
            suffix = ' (cont.)' if page_no else ''
            self._slide_title(slide, oslide.title + suffix,
                              subhead=oslide.subhead if page_no == 0 else '',
                              kicker=oslide.kicker)
            self._render_blocks(slide, blocks, font, body_y)
            if oslide.notes:
                notes_tf = slide.notes_slide.notes_text_frame
                text = oslide.notes if page_no == 0 \
                    else 'See notes on the previous slide.'
                first = True
                for para in text.split('\n\n'):
                    p = _para(notes_tf, first)
                    first = False
                    p.text = para

    def _render_blocks(self, slide, blocks: list[OBlock], font: float,
                       y_start: float = BODY_Y) -> None:
        y = y_start
        text_group: list[OBlock] = []

        def flush_text() -> None:
            nonlocal y, text_group
            if not text_group:
                return
            height = sum(self._block_height(b, font) for b in text_group)
            _, tf = _textbox(slide, BODY_X, y, BODY_W, min(height + 0.1,
                                                           7.0 - y))
            first = True
            for b in text_group:
                p = _para(tf, first)
                first = False
                p.space_after = Pt(5)
                if b.kind == 'para':
                    r = p.add_run()
                    r.text = b.text
                    _set_font(r, font, INK)
                elif b.kind == 'subheading':
                    r = p.add_run()
                    r.text = b.text
                    _set_font(r, 18 if b.level == 3 else 15, NAVY, bold=True)
                    p.space_before = Pt(8)
                else:  # bullet
                    lvl = min(b.level, 4)
                    _set_indent(p, 0.30 * lvl + 0.18, hang_in=0.18)
                    glyph = '▪' if b.level <= 1 else '–'
                    r = p.add_run()
                    r.text = glyph + '  '
                    _set_font(r, font, NAVY if b.level <= 1 else MUTED)
                    r = p.add_run()
                    r.text = b.text
                    _set_font(r, font, INK)
            y += height + 0.08
            text_group = []

        for b in blocks:
            if b.kind in ('para', 'bullet', 'subheading'):
                text_group.append(b)
                continue
            flush_text()
            if b.kind == 'table':
                y += self._render_table(slide, b, BODY_X, y, BODY_W) + 0.15
            elif b.kind == 'image':
                y += self._render_image(slide, b, y) + 0.12
            elif b.kind == 'persona':
                y += self._render_persona(slide, b, y, font) + 0.15
        flush_text()

    def _render_persona(self, slide, b: OBlock, y: float,
                        font: float) -> float:
        text_h = _text_height(b.text, font, BODY_W - 1.1)
        card_h = text_h + 0.42
        _rounded_rect(slide, BODY_X, y, BODY_W, card_h, NAVY_2, radius=0.12)
        _, tf = _textbox(slide, BODY_X + 0.32, y, 0.6, card_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = '“'
        _set_font(r, 40, ICE, bold=True, font=FONT_HEAD)
        _, tf = _textbox(slide, BODY_X + 0.95, y, BODY_W - 1.25, card_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = b.text
        _set_font(r, font, WHITE, italic=True)
        return card_h

    def _render_table(self, slide, b: OBlock, x, y, max_w) -> float:
        rows = b.rows
        n_rows, n_cols = len(rows), max(len(r) for r in rows)
        widths = _table_col_widths(rows, max_w)
        table_w = sum(widths)
        font = 12
        if any(len(cell) > 120 for r in rows for cell in r):
            font = 10
        heights = _table_row_heights(rows, widths, font)
        gfx = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                     Inches(table_w),
                                     Inches(sum(heights)))
        table = gfx.table
        table.first_row = False   # we style manually
        table.horz_banding = False
        for c, w in enumerate(widths):
            table.columns[c].width = Inches(w)
        for ri in range(n_rows):
            table.rows[ri].height = Inches(heights[ri])
            for ci in range(n_cols):
                cell = table.cell(ri, ci)
                cell.fill.solid()
                if ri == 0:
                    cell.fill.fore_color.rgb = NAVY
                else:
                    cell.fill.fore_color.rgb = WHITE if ri % 2 else ICE_BG
                cell.margin_left = cell.margin_right = Inches(0.08)
                cell.margin_top = cell.margin_bottom = Inches(0.03)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = rows[ri][ci] if ci < len(rows[ri]) else ''
                _set_font(r, font, WHITE if ri == 0 else INK, bold=ri == 0)
        return sum(heights)

    def _render_image(self, slide, b: OBlock, y: float) -> float:
        path = self._resolve_media(b.path)
        if path:
            try:
                from PIL import Image as PILImage
                with PILImage.open(path) as im:
                    w_px, h_px = im.size
                max_w, max_h = 7.5, 3.2
                w_in, h_in = w_px / 96, h_px / 96
                scale = min(max_w / w_in, max_h / h_in, 1.0)
                w_in, h_in = w_in * scale, max(0.6, h_in * scale)
                slide.shapes.add_picture(path, Inches(BODY_X), Inches(y),
                                         Inches(w_in), Inches(h_in))
                return h_in
            except Exception:
                pass  # unreadable / unsupported raster -> placeholder
        # styled placeholder — a missing image must never fail the build
        ph_h = 1.4
        _rounded_rect(slide, BODY_X, y, 7.5, ph_h, ICE_BG, radius=0.1)
        _, tf = _textbox(slide, BODY_X + 0.3, y, 6.9, ph_h,
                         anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f'Image: {os.path.basename(b.path) or "(unnamed)"}'
        _set_font(r, 13, MUTED, bold=True)
        p = tf.add_paragraph()
        where = self.media_dir or os.path.join(self.base_dir,
                                               os.path.dirname(b.path) or '.')
        r = p.add_run()
        r.text = f'not found under {where}'
        _set_font(r, 11, MUTED)
        if b.alt and b.alt != os.path.basename(b.path):
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = b.alt
            _set_font(r, 11, MUTED, italic=True)
        return ph_h

    # -- takeaways slide -------------------------------------------------
    def _takeaways_slide(self, oslide: OSlide) -> None:
        slide = self._new_slide(dark=True)
        self._slide_title(slide, oslide.title or 'Key takeaways', dark=True)
        items = oslide.items
        row_h = min(1.15, BODY_H / max(1, len(items)))
        y = BODY_Y + 0.15
        for idx, item in enumerate(items):
            badge = _rounded_rect(slide, BODY_X, y, 0.5, 0.5, ICE, radius=0.5)
            btf = badge.text_frame
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            btf.word_wrap = False
            p = btf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(idx + 1)
            _set_font(r, 16, NAVY, bold=True)
            size = 19 if _text_lines(item, 19, BODY_W - 0.9) <= 2 else 16
            _, tf = _textbox(slide, BODY_X + 0.75, y - 0.04, BODY_W - 0.9,
                             row_h)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = item
            _set_font(r, size, ICE)
            y += row_h

    # -- related-documents closer ------------------------------------------
    def _related_slides(self) -> None:
        entries = self.o.related
        per_slide = 5
        for start in range(0, len(entries), per_slide):
            chunk = entries[start:start + per_slide]
            slide = self._new_slide(dark=True)
            suffix = ' (cont.)' if start else ''
            self._slide_title(slide, 'Related documents' + suffix, dark=True)
            y = BODY_Y + 0.05
            card_h = min(1.02, (BODY_H - 0.1) / len(chunk) - 0.12)
            for e in chunk:
                _rounded_rect(slide, BODY_X, y, BODY_W, card_h, NAVY_2,
                              radius=0.1)
                title = e.get('title') or e.get('why') or 'Related document'
                _, tf = _textbox(slide, BODY_X + 0.3, y + 0.1,
                                 BODY_W - 1.6, card_h - 0.2)
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = title
                _set_font(r, 16, ICE, bold=True)
                why = e.get('why') if e.get('title') else ''
                if why:
                    p = tf.add_paragraph()
                    r = p.add_run()
                    r.text = why
                    _set_font(r, 11.5, LIGHT_ICE)
                    p.space_before = Pt(3)
                if e.get('doc') is not None:
                    _chip(slide, SLIDE_W_IN - MARGIN - 1.1, y + 0.12,
                          f"doc {e['doc']}", fill=ICE, color=NAVY, size=10,
                          bold=True, h=0.3)
                y += card_h + 0.14
            strip = ' · '.join(x for x in (
                self.o.source_file,
                f"Extracted {_fmt_date(_meta_str(self.o.meta, 'extracted'))}"
                if _meta_str(self.o.meta, 'extracted') else '') if x)
            if strip:
                _, tf = _textbox(slide, MARGIN, 7.02, BODY_W, 0.35)
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = strip
                _set_font(r, 10, ICE_MUTED)


# ------------------------------------------------------------------- API

def convert(md_text: str, media_dir: str | None = None,
            base_dir: str = '.', enhance: bool = False,
            model: str = 'claude-opus-5',
            strict: bool = False) -> tuple[Presentation, DeckOutline, str]:
    """In-process entry point. Returns (presentation, outline, mode)
    where mode is 'enhanced' or 'deterministic'."""
    sc = parse_sidecar(md_text)
    outline, mode = outline_from_sidecar(sc), 'deterministic'
    if enhance:
        from enhance import EnhanceError, enhance_outline
        try:
            outline, mode = enhance_outline(sc, md_text, model=model), 'enhanced'
        except EnhanceError as exc:
            if strict:
                raise
            print(f'warning: enhancement unavailable ({exc}); '
                  f'falling back to deterministic outline', file=sys.stderr)
    prs = DeckRenderer(outline, media_dir, base_dir).build()
    return prs, outline, mode


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(
        description='Render an LRS Doc Index sidecar .md as a .pptx '
                    'briefing deck.')
    ap.add_argument('input', help='sidecar markdown file')
    ap.add_argument('-o', '--output', help='output .pptx path '
                    '(default: input stem + .pptx)')
    ap.add_argument('--media-dir', help='directory holding the images the '
                    'sidecar links (default: resolve links relative to the '
                    'input file)')
    ap.add_argument('--enhance', action='store_true',
                    help='restructure the deck with the Claude API '
                    '(pip install anthropic + ANTHROPIC_API_KEY); '
                    'falls back to the deterministic outline on failure')
    ap.add_argument('--model', default='claude-opus-5',
                    help='Claude model for --enhance (default: %(default)s)')
    ap.add_argument('--strict-enhance', action='store_true',
                    help='exit 3 instead of falling back when --enhance fails')
    ap.add_argument('--outline-json', metavar='PATH',
                    help='write the outline used to PATH (audit artifact); '
                    'with --from-outline, read the outline from PATH instead')
    ap.add_argument('--from-outline', action='store_true',
                    help='render --outline-json PATH instead of building an '
                    'outline (metadata still comes from the input sidecar)')
    args = ap.parse_args(argv)

    try:
        md_text = open(args.input, encoding='utf-8').read()
    except OSError as exc:
        print(f'error: cannot read {args.input}: {exc}', file=sys.stderr)
        return 2

    base_dir = os.path.dirname(os.path.abspath(args.input))
    out = args.output or os.path.splitext(args.input)[0] + '.pptx'

    try:
        if args.from_outline:
            if not args.outline_json:
                print('error: --from-outline requires --outline-json PATH',
                      file=sys.stderr)
                return 2
            sc = parse_sidecar(md_text)
            data = json.load(open(args.outline_json, encoding='utf-8'))
            outline = load_outline(data, sc)
            prs = DeckRenderer(outline, args.media_dir, base_dir).build()
            mode = 'from-outline'
        else:
            try:
                prs, outline, mode = convert(
                    md_text, media_dir=args.media_dir, base_dir=base_dir,
                    enhance=args.enhance, model=args.model,
                    strict=args.strict_enhance)
            except Exception as exc:
                if type(exc).__name__ == 'EnhanceError':
                    print(f'error: enhancement failed under '
                          f'--strict-enhance: {exc}', file=sys.stderr)
                    return 3
                raise
            if args.outline_json:
                with open(args.outline_json, 'w', encoding='utf-8') as fh:
                    json.dump(outline.to_dict(), fh, indent=2)
    except SidecarError as exc:
        print(f'error: {args.input} is not a sidecar: {exc}', file=sys.stderr)
        return 2

    prs.save(out)
    print(f'{out}: {len(prs.slides._sldIdLst)} slides ({mode})')
    return 0


if __name__ == '__main__':
    sys.exit(main())
